#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
CDMS NAS Worker
================
лѓЄл¶ЉлѓЄлФФмЦілЮ© CDMSмЭШ NAS мЧ∞лПЩ мЫМмї§.
- SupabaseмЭШ nas_tasks нБРл•Љ нПілІБнХШмЧђ мЮСмЧЕмЭД м≤Шл¶ђнХШк≥†,
- м£ЉкЄ∞м†БмЬЉл°Ь м†Дм≤і к≥Љм†Х нПілНФл•Љ мК§мЇФнХШмЧђ м†ЬмЮСлЛ®к≥Д мІДнЦЙмГБнГЬ(нММмЭЉ мЬ†лђі¬ЈмИШм†ХмЭЉ)мЩА
  мҐЕнОЄ мШБмГБкЄЄмЭі(ffprobe)л•Љ мЮРлПЩмЬЉл°Ь м±ДмЪілЛ§.

м≤Шл¶ђ action:
  ping            : мЧ∞к≤∞ нЩХмЭЄ / к≥µмЬ†нПілНФ л™©л°Э
  mkdir_tree      : к≥Љм†Х нПілНФ + лЛ®к≥Д нХШмЬДнПілНФ мГЭмД± (params: project, folders)
  scan_progress   : к≥Љм†Х лЛ®к≥ДнПілНФ мК§мЇФ вЖТ lesson_stage(status,file_mtime,file_name) к∞±мЛ†
                    + мҐЕнОЄ нММмЭЉ ffprobe вЖТ lessons.duration_sec к∞±мЛ†
  probe_durations : мҐЕнОЄ мШБмГБкЄЄмЭілІМ лЛ§мЛЬ мґФмґЬ
  rename_folder   : NAS нПілНФл™Е л≥Ак≤љ (params: old, new)  [CDMSвЖТNAS]
  sync_names      : NAS нПілНФл™ЕмЭД мЭљмЦі CDMSмЧР л∞ШмШБ           [NASвЖТCDMS] (кЄ∞міИ кµђнШД)

вЪ†пЄП мЭі мЫМмї§лКФ Supabase SERVICE ROLE нВ§л•Љ мВђмЪ©нХЬлЛ§. м†ИлМА лЄМлЭЉмЪ∞м†А/нФДлЯ∞нКЄ(index.html)мЧР
   лД£мІА лІР к≤Г. мДЬл≤Д(.env)мЧРлІМ л≥ікіАнХЬлЛ§.
"""
import os, re, sys, time, json, tempfile, subprocess, traceback
from datetime import datetime, timezone, timedelta

try:
    from supabase import create_client, Client
except ImportError:
    print("supabase нМ®нВ§мІАк∞А нХДмЪФнХ©лЛИлЛ§:  pip install supabase", file=sys.stderr)
    raise

# ----------------------------------------------------------------------------
# мД§м†Х (.env / нЩШк≤љл≥АмИШ)
# ----------------------------------------------------------------------------
SB_URL   = os.environ.get("SUPABASE_URL", "https://kowtvvrgpzgrdlnxasxw.supabase.co")
SB_KEY   = os.environ.get("SUPABASE_SERVICE_ROLE_KEY")            # нХДмИШ (service_role)
NAS_MODE = os.environ.get("NAS_MODE", "mount").lower()           # mount | smb

# mount л™®лУЬ: NASк∞А мДЬл≤ДмЧР лІИмЪінКЄлРЬ к≤љл°Ь
NAS_BASE = os.environ.get("NAS_BASE", "/mnt/nas")

# smb л™®лУЬ: nas_config нЕМмЭілЄФ к∞ТмЭД кЄ∞л≥ЄмЬЉл°Ь мУ∞лРШ нЩШк≤љл≥АмИШл°Ь лНЃмЦімУЄ мИШ мЮИмЭМ
SMB_HOST = os.environ.get("NAS_SMB_HOST", "")
SMB_USER = os.environ.get("NAS_SMB_USER", "")
SMB_PASS = os.environ.get("NAS_SMB_PASS", "")
SMB_SHARE= os.environ.get("NAS_SMB_SHARE", "")

# кЄЄмЭі мґФмґЬ кЄ∞м§А лЛ®к≥Д (мҐЕнОЄ=7). нПіл∞± мЧЖмЭМ(мВђмЪ©мЮР мД†нГЭ: мҐЕнОЄ кЄ∞м§А)
LENGTH_STAGE_ID = int(os.environ.get("LENGTH_STAGE_ID", "7"))
# мШБмГБ нЩХмЮ•мЮР
VIDEO_EXT = {".mp4", ".mov", ".mkv", ".avi", ".m4v", ".wmv", ".mpg", ".mpeg", ".mts"}
# мЮРлПЩ мК§мЇФ м£ЉкЄ∞(міИ). 0 мЭіл©і нБР мЮСмЧЕлІМ м≤Шл¶ђ(м£ЉкЄ∞ мК§мЇФ лБФ)
SCAN_INTERVAL = int(os.environ.get("SCAN_INTERVAL", "600"))
POLL_INTERVAL = int(os.environ.get("POLL_INTERVAL", "3"))

# мЮРлПЩ мЭіл©ФмЭЉ (лЛ®к≥Д мЩДл£М вЖТ лЛ§мЭМ лЛілЛємЮР)
EMAIL_ENABLED  = os.environ.get("EMAIL_ENABLED", "true").lower() == "true"
EMAIL_PROVIDER = os.environ.get("EMAIL_PROVIDER", "resend").lower()   # resend | sendgrid | smtp
EMAIL_API_KEY  = os.environ.get("EMAIL_API_KEY", "")
EMAIL_FROM     = os.environ.get("EMAIL_FROM", "CDMS <noreply@mirimmedialab.co.kr>")
CDMS_URL       = os.environ.get("CDMS_URL", "https://cdms.mirimmedialab.co.kr")
# SMTP (EMAIL_PROVIDER=smtp мЭЉ лХМ вАФ мШИ: нХШмЭімЫНмК§ smtps.hiworks.com:465 SSL)
SMTP_HOST      = os.environ.get("SMTP_HOST", "smtps.hiworks.com")
SMTP_PORT      = int(os.environ.get("SMTP_PORT", "465"))
SMTP_USER      = os.environ.get("SMTP_USER", "")
SMTP_PASS      = os.environ.get("SMTP_PASS", "")
# мХМл¶Љ л©ФмЭЉ нПіл∞± мИШмЛ†мЮР(PM мЭіл©ФмЭЉмЭі мЮДмЛЬ/лѓЄмД§м†ХмЭЉ лХМ)
REMIND_EMAIL   = os.environ.get("REMIND_EMAIL", "ghnam7312@gmail.com")
# мШБмГБк≤АмИШ: мЛЬлЖАл°ЬмІА лУ±мЬЉл°Ь NAS мШБмГБмЭД мЩЄлґА HTTPS мІБм†С мДЬлєЩнХШлКФ к≤љмЪ∞.
#  NAS_PUBLIC_BASE к∞А мД§м†ХлРШл©і Supabase мЧЕл°ЬлУЬ лМАмЛ† к≥µк∞Ь URLмЭД к≤АмИШмШБмГБмЬЉл°Ь мВђмЪ©нХЬлЛ§.
#  мШИ) NAS_PUBLIC_ROOT=/mnt  (Web Station лђЄмДЬл£®нКЄк∞А /mnt мЧР лІ§нХС)
#      NAS_PUBLIC_BASE=https://nas.mirimmedialab.co.kr/files  (кЈЄ л£®нКЄмЭШ к≥µк∞Ь л≤†мЭімК§)
NAS_PUBLIC_BASE = os.environ.get("NAS_PUBLIC_BASE", "")
NAS_PUBLIC_ROOT = os.environ.get("NAS_PUBLIC_ROOT", NAS_BASE)

# м∞Єм°∞мЛЬнКЄ кЄ∞м§А лЛ®к≥Д мИЬмДЬ/лЭЉл≤® (нФДлЯ∞нКЄмЩА лПЩмЭЉ) вАФ "лЛ§мЭМ лЛ®к≥Д" к≥ДмВ∞мЪ©
REF_ORDER = [1, 2, 3, 4, 5, 6, 7, 10, 13, 9]
REF_LABEL = {1: "мЫРк≥†", 2: "міђмШБ", 3: "к∞АнОЄ", 4: "мК§нБђл¶љнКЄ", 5: "мК§нЖ†л¶ђл≥ілУЬ",
             6: "лФФмЮРмЭЄ", 7: "мҐЕнОЄ", 10: "srt", 13: "л≤ИмЧ≠", 9: "нХЩмКµмЮРл£М",
             8: "к≤АмИШ", 11: "мЭМмД±", 12: "HTML", 0: "лђЄмДЬ"}


def stage_key(sid):
    return REF_ORDER.index(sid) if sid in REF_ORDER else 100 + sid

if not SB_KEY:
    print("нЩШк≤љл≥АмИШ SUPABASE_SERVICE_ROLE_KEY к∞А нХДмЪФнХ©лЛИлЛ§.", file=sys.stderr)
    sys.exit(1)

sb: Client = create_client(SB_URL, SB_KEY)

# м∞®мЛЬ л≤ИнШЄ мґФмґЬ:  "3м∞®мЛЬ", "03м∞®мЛЬ", "3 м∞®мЛЬ"  вЖТ  3
RE_LESSON = re.compile(r"(\d+)\s*м∞®\s*мЛЬ")
# м£Љм∞® л≤ИнШЄ(мЮИмЬЉл©і нХ®кїШ лІ§мє≠):  "7м£Љм∞®"
RE_WEEK   = re.compile(r"(\d+)\s*м£Љ\s*м∞®")


def log(*a):
    print(datetime.now().strftime("%H:%M:%S"), *a, flush=True)


def _from_addr(s):
    m = re.search(r"<([^>]+)>", s or "")
    return m.group(1) if m else (s or "")


def _send_smtp(to_email, subject, html):
    """SMTP(SSL) л∞ЬмЖ° вАФ нХШмЭімЫНмК§ лУ± кЄ∞м°і л©ФмЭЉмДЬл≤ДмЪ©."""
    import smtplib, ssl
    from email.mime.text import MIMEText
    from email.utils import formataddr, parseaddr
    if not (SMTP_USER and SMTP_PASS):
        return (False, "SMTP_USER/SMTP_PASS мЧЖмЭМ")
    name, addr = parseaddr(EMAIL_FROM)
    from_addr = addr or SMTP_USER
    msg = MIMEText(html, "html", "utf-8")
    msg["Subject"] = subject
    msg["From"] = formataddr((name or "CDMS", from_addr))
    msg["To"] = to_email
    try:
        ctx = ssl.create_default_context()
        with smtplib.SMTP_SSL(SMTP_HOST, SMTP_PORT, context=ctx, timeout=20) as s:
            s.login(SMTP_USER, SMTP_PASS)
            s.sendmail(from_addr, [to_email], msg.as_string())
        return (True, None)
    except Exception as e:
        return (False, str(e)[:200])


def send_email(to_email, subject, html):
    """Resend / SendGrid HTTP API лШРлКФ SMTP л°Ь л©ФмЭЉ л∞ЬмЖ°. (True, None) лШРлКФ (False, error)."""
    if not (EMAIL_ENABLED and to_email):
        return (False, "email лєДнЩЬмД±/мИШмЛ†мЮРмЧЖмЭМ")
    if EMAIL_PROVIDER == "smtp":
        return _send_smtp(to_email, subject, html)
    import urllib.request, urllib.error
    if not EMAIL_API_KEY:
        return (False, "EMAIL_API_KEY мЧЖмЭМ")
    try:
        if EMAIL_PROVIDER == "sendgrid":
            url = "https://api.sendgrid.com/v3/mail/send"
            payload = {"personalizations": [{"to": [{"email": to_email}]}],
                       "from": {"email": _from_addr(EMAIL_FROM)},
                       "subject": subject,
                       "content": [{"type": "text/html", "value": html}]}
        else:  # resend
            url = "https://api.resend.com/emails"
            payload = {"from": EMAIL_FROM, "to": [to_email],
                       "subject": subject, "html": html}
        req = urllib.request.Request(
            url, data=json.dumps(payload).encode("utf-8"),
            headers={"Authorization": "Bearer %s" % EMAIL_API_KEY,
                     "Content-Type": "application/json",
                     "Accept": "application/json",
                     # Cloudflareк∞А кЄ∞л≥Є python-urllib UAл•Љ ліЗ(1010)мЬЉл°Ь м∞®лЛ®нХШлѓАл°Ь мЭЉл∞Ш UA мІАм†Х
                     "User-Agent": "Mozilla/5.0 (X11; Linux x86_64) CDMS-NAS-Worker/1.0"})
        urllib.request.urlopen(req, timeout=20).read()
        return (True, None)
    except urllib.error.HTTPError as e:
        try:
            body = e.read().decode()[:200]
        except Exception:
            body = ""
        return (False, "HTTP %s %s" % (e.code, body))
    except Exception as e:
        return (False, str(e))


# ============================================================================
# нММмЭЉмЛЬмК§нЕЬ мґФмГБнЩФ (mount / smb)
# ============================================================================
class MountFS:
    """NASк∞А л°ЬмїђмЧР лІИмЪінКЄлРЬ к≤љмЪ∞. мЭЉл∞Ш нММмЭЉмЛЬмК§нЕЬ мВђмЪ©."""
    def __init__(self, base):
        self.base = base

    def _abs(self, path):
        if not path:
            return self.base
        if os.path.isabs(path):
            return path
        return os.path.join(self.base, path)

    def exists(self, path):
        return os.path.exists(self._abs(path))

    def listfiles(self, path):
        p = self._abs(path)
        if not os.path.isdir(p):
            return []
        out = []
        for name in os.listdir(p):
            fp = os.path.join(p, name)
            if os.path.isfile(fp):
                out.append((name, os.path.getmtime(fp)))
        return out

    def walkfiles(self, path):
        """лЛ®к≥Д нПілНФл•Љ нХШмЬДнПілНФкєМмІА нЫСмЦі (лЛ®к≥ДнПілНФ кЄ∞м§А мГБлМАк≤љл°Ь, mtime) л™©л°Э л∞ШнЩШ."""
        base = self._abs(path)
        if not os.path.isdir(base):
            return []
        out = []
        for root, dirs, files in os.walk(base):
            dirs[:] = [d for d in dirs if not d.startswith(".")]
            for f in files:
                if f.startswith(".") or f == "Thumbs.db":
                    continue
                fp = os.path.join(root, f)
                try:
                    mt = os.path.getmtime(fp)
                except OSError:
                    continue
                out.append((os.path.relpath(fp, base).replace("\\", "/"), mt))
        return out

    def listdirs(self, path):
        p = self._abs(path)
        if not os.path.isdir(p):
            return []
        return [n for n in os.listdir(p) if os.path.isdir(os.path.join(p, n))]

    def makedirs(self, path):
        os.makedirs(self._abs(path), exist_ok=True)

    def rename(self, old, new):
        os.rename(self._abs(old), self._abs(new))

    def local_copy(self, path):
        """ffprobeмЪ© л°Ьмїђ к≤љл°Ь. лІИмЪінКЄ л™®лУЬлКФ кЈЄлМАл°Ь л∞ШнЩШ."""
        return self._abs(path), False  # (к≤љл°Ь, мЮДмЛЬнММмЭЉмЧђлґА)

    def shares(self):
        try:
            return [os.path.join(self.base, n) for n in os.listdir(self.base)]
        except Exception as e:
            return ["(л™©л°Э мЛ§нМ®: %s)" % e]


class SmbFS:
    """SMB мІБм†С м†СмЖН. smbprotocol(smbclient) мВђмЪ©."""
    def __init__(self, host, user, pw, share):
        import smbclient
        self.smbclient = smbclient
        self.host = host
        self.share = share
        smbclient.ClientConfig(username=user, password=pw)
        self._user = user
        self._pw = pw

    def _unc(self, path):
        if path and path.startswith("\\\\"):   # мЭілѓЄ м†Дм≤і UNC к≤љл°Ьл©і кЈЄлМАл°Ь
            return path.replace("/", "\\")
        path = (path or "").replace("/", "\\").lstrip("\\")
        base = r"\\%s\%s" % (self.host, self.share)
        return base + ("\\" + path if path else "")

    def exists(self, path):
        try:
            self.smbclient.stat(self._unc(path)); return True
        except Exception:
            return False

    def listfiles(self, path):
        out = []
        try:
            for e in self.smbclient.scandir(self._unc(path)):
                if e.is_file():
                    out.append((e.name, e.stat().st_mtime))
        except Exception:
            pass
        return out

    def walkfiles(self, path):
        out = []
        def rec(rel):
            sub = path + ("/" + rel if rel else "")
            try:
                entries = list(self.smbclient.scandir(self._unc(sub)))
            except Exception:
                return
            for e in entries:
                nm = e.name
                if nm.startswith(".") or nm == "Thumbs.db":
                    continue
                childrel = (rel + "/" + nm) if rel else nm
                try:
                    if e.is_dir():
                        rec(childrel)
                    elif e.is_file():
                        out.append((childrel, e.stat().st_mtime))
                except Exception:
                    continue
        rec("")
        return out

    def listdirs(self, path):
        out = []
        try:
            for e in self.smbclient.scandir(self._unc(path)):
                if e.is_dir():
                    out.append(e.name)
        except Exception:
            pass
        return out

    def makedirs(self, path):
        parts, cur = path.replace("\\", "/").strip("/").split("/"), ""
        for p in parts:
            cur = (cur + "/" + p) if cur else p
            try:
                self.smbclient.mkdir(self._unc(cur))
            except Exception:
                pass

    def rename(self, old, new):
        self.smbclient.rename(self._unc(old), self._unc(new))

    def local_copy(self, path):
        tmp = tempfile.NamedTemporaryFile(delete=False, suffix=os.path.splitext(path)[1])
        with self.smbclient.open_file(self._unc(path), mode="rb") as f:
            tmp.write(f.read())
        tmp.close()
        return tmp.name, True

    def shares(self):
        return [r"\\%s\%s" % (self.host, self.share)]


def build_fs():
    if NAS_MODE == "smb":
        host, user, pw, share = SMB_HOST, SMB_USER, SMB_PASS, SMB_SHARE
        # лєДмЦі мЮИмЬЉл©і nas_configмЧРмДЬ л≥імґ©
        if not (host and user and share):
            try:
                cfg = sb.table("nas_config").select("*").limit(1).execute().data
                if cfg:
                    c = cfg[0]
                    host = host or (c.get("url") or "").replace("smb://", "").strip("/")
                    user = user or c.get("username") or ""
                    pw   = pw   or c.get("password") or ""
                    share= share or c.get("base") or ""
            except Exception as e:
                log("nas_config мЭљкЄ∞ мЛ§нМ®:", e)
        return SmbFS(host, user, pw, share)
    return MountFS(NAS_BASE)


# ============================================================================
# ffprobe
# ============================================================================
def ffprobe_seconds(fs, remote_path):
    local, is_tmp = fs.local_copy(remote_path)
    try:
        out = subprocess.run(
            ["ffprobe", "-v", "error", "-show_entries", "format=duration",
             "-of", "default=nw=1:nokey=1", local],
            capture_output=True, text=True, timeout=120)
        val = (out.stdout or "").strip()
        return int(round(float(val))) if val else None
    except Exception as e:
        log("ffprobe мЛ§нМ®:", remote_path, e)
        return None
    finally:
        if is_tmp:
            try: os.unlink(local)
            except Exception: pass


# ============================================================================
# лН∞мЭінД∞ л°ЬлФ©
# ============================================================================
def load_project_bundle(project_id):
    proj = sb.table("projects").select("*").eq("id", project_id).limit(1).execute().data
    if not proj:
        return None
    proj = proj[0]
    stages = {s["id"]: s for s in sb.table("stages").select("*").execute().data}
    pstages = sb.table("project_stages").select("*").eq("project_id", project_id).eq("enabled", True).execute().data
    enabled = [stages[p["stage_id"]] for p in pstages if p["stage_id"] in stages]
    lessons = sb.table("lessons").select("id,lesson_no,week_id").eq("project_id", project_id).execute().data
    weeks = {w["id"]: w["week_no"] for w in sb.table("weeks").select("id,week_no").eq("project_id", project_id).execute().data}
    for l in lessons:
        l["week_no"] = weeks.get(l.get("week_id"))
    lesson_ids = [l["id"] for l in lessons] or ["00000000-0000-0000-0000-000000000000"]
    ls = sb.table("lesson_stage").select("lesson_id,stage_id,status,file_name").in_(
        "lesson_id", lesson_ids).execute().data
    cur = {(r["lesson_id"], r["stage_id"]): r for r in ls}
    # лЛ®к≥Дл≥Д лЛілЛємЮР + мВђмЪ©мЮР + кЄ∞м°і л∞ЬмЖ°л°ЬкЈЄ(м§Сл≥µл∞©мІА)
    assignees = {a["stage_id"]: a["user_id"]
                 for a in sb.table("stage_assignees").select("stage_id,user_id").eq("project_id", project_id).execute().data}
    users = {u["id"]: u for u in sb.table("users").select("id,name,email").execute().data}
    notified = set()
    for r in sb.table("email_notifications").select("lesson_id,completed_stage_id").in_("lesson_id", lesson_ids).execute().data:
        notified.add((r["lesson_id"], r["completed_stage_id"]))
    return {"proj": proj, "enabled": enabled, "lessons": lessons, "ls": cur,
            "assignees": assignees, "users": users, "notified": notified}


def match_lesson(filename, lessons, has_weeks):
    m = RE_LESSON.search(filename)
    wk = RE_WEEK.search(filename)
    lesson_no = int(m.group(1)) if m else None
    week_no = int(wk.group(1)) if wk else None
    if has_weeks and wk and m:
        for l in lessons:
            if l["lesson_no"] == lesson_no and l.get("week_no") == week_no:
                return l
    if (not has_weeks) and week_no is not None:
        # м∞®мЛЬнШХ к≥Љм†Х: нММмЭЉл™ЕмЭШ "Nм£Љм∞®"к∞А CDMS м∞®мЛЬ л≤ИнШЄ, нММмЭЉл™ЕмЭШ "Nм∞®мЛЬ"лКФ нММнКЄ л≤ИнШЄ
        #  (мШИ: лН∞мЭінД∞мВђмЭімЦЄмК§_01м£Љм∞® 2м∞®мЛЬ.mp4 = 1м∞®мЛЬмЭШ 2л≤ИмІЄ нММнКЄ) вАФ Edge(nas-versions matchLesson)мЩА лПЩмЭЉ кЈЬмєЩ (2026-07-17)
        cands = [l for l in lessons if l["lesson_no"] == week_no]
        if len(cands) == 1:
            return cands[0]
    if lesson_no is None:
        return None
    cands = [l for l in lessons if l["lesson_no"] == lesson_no]
    if len(cands) == 1:
        return cands[0]
    return None


# ============================================================================
# action: scan_progress
# ============================================================================
def action_scan_progress(fs, project_id, do_duration=False):
    # (2026-07-15) do_duration кЄ∞л≥Є False: мШБмГБ кЄЄмЭілКФ Edge(nas-versions)к∞А нММнКЄ "нХ©мВ∞"мЬЉл°Ь к≥ДмВ∞¬ЈкіАл¶ђ.
    # мЫМмї§мЭШ лЛ®мЭЉ нММмЭЉ(мµЬкЈЉ 1к∞Ь) кЄЄмЭі лНЃмЦімУ∞кЄ∞лКФ лЛ§м§С нММнКЄ м∞®мЛЬмЭШ нХ©к≥Дл•Љ нЫЉмЖРнХШлѓАл°Ь лєДнЩЬмД±нЩФ.
    b = load_project_bundle(project_id)
    if not b:
        return {"ok": False, "error": "project not found"}
    proj, enabled, lessons, cur = b["proj"], b["enabled"], b["lessons"], b["ls"]
    assignees, users, notified = b["assignees"], b["users"], b["notified"]
    root = proj.get("nas_root")
    if not root:
        return {"ok": False, "error": "nas_root лѓЄмД§м†Х вАФ л®Љм†А NAS нПілНФмГЭмД±(mkdir_tree)мЭД мЛ§нЦЙнХШмДЄмЪФ"}
    if re.match(r"^nas\d+:", root):
        return {"ok": True, "skipped": True, "note": "лЛ§л•Є NAS(nasN:) к≥Љм†Х вАФ Edge(nas-versions) мК§мЇФмЭі лЛілЛє"}
    has_weeks = any(l.get("week_no") for l in lessons)
    now = datetime.now(timezone.utc).isoformat()

    found = {}   # (lesson_id, stage_id) -> (mtime, name)
    durations = {}  # lesson_id -> remote path

    for s in enabled:
        folder = "%s/%s" % (root, s["nas_folder"])
        if not fs.exists(folder):
            continue
        for rel, mtime in fs.walkfiles(folder):
            base = rel.rsplit("/", 1)[-1]
            if base.startswith("~") or ".cdms_" in rel:
                continue
            # нММмЭЉл™Е лШРлКФ мГБмЬД м∞®мЛЬнПілНФл™Е(мШИ: 07_мҐЕнОЄ/2м∞®мЛЬ/мШБмГБ.mp4)мЬЉл°Ь лІ§мє≠
            l = match_lesson(rel, lessons, has_weeks)
            if not l:
                continue
            key = (l["id"], s["id"])
            if key not in found or mtime > found[key][0]:
                found[key] = (mtime, base)
            # кЄЄмЭі: мҐЕнОЄ(LENGTH_STAGE_ID) мШБмГБнММмЭЉ вАФ к∞ЩмЭА м∞®мЛЬмЧР мЧђлЯњмЭіл©і к∞АмЮ• мµЬкЈЉ
            if do_duration and s["id"] == LENGTH_STAGE_ID and os.path.splitext(base)[1].lower() in VIDEO_EXT:
                if l["id"] not in durations or mtime >= durations[l["id"]][0]:
                    durations[l["id"]] = (mtime, "%s/%s" % (folder, rel))

    upserts, reverts, updated_dur = [], [], 0

    # 1) нММмЭЉ л∞Ьк≤ђ вЖТ done + мИШм†ХмЭЉ + нММмЭЉл™Е
    for (lesson_id, stage_id), (mtime, name) in found.items():
        iso = datetime.fromtimestamp(mtime, tz=timezone.utc).isoformat()
        upserts.append({"lesson_id": lesson_id, "stage_id": stage_id, "status": "done",
                        "file_mtime": iso, "file_name": name, "updated_at": now})

    # 2) мЭім†ДмЧР мЮРлПЩнСЬкЄ∞(file_name мЮИмЭМ)мШАмЬЉлВШ мІАкЄИ нММмЭЉмЭі мВђлЭЉмІД к≤љмЪ∞ вЖТ wait л°Ь лРШлПМл¶Љ
    for (lesson_id, stage_id), row in cur.items():
        if (lesson_id, stage_id) in found:
            continue
        if row.get("file_name"):  # мЮРлПЩмЬЉл°Ь м±ДмЫМм°МлНШ мЕА
            reverts.append({"lesson_id": lesson_id, "stage_id": stage_id, "status": "wait",
                            "file_mtime": None, "file_name": None, "updated_at": now})

    # (2026-07) мІДнЦЙ нСЬкЄ∞/лРШлПМл¶ЉмЭА Edge(nas-versions scan)к∞А нММмЭЉ "мГЭмД±мЭЉ"¬ЈмИШм†Хл≥Є¬Јк≤АмИШмВђмЭінБі кЄ∞м§АмЬЉл°Ь лЛілЛє.
    # мЫМмї§лКФ мШБмГБкЄЄмЭі(ffprobe)мЩА мЩДл£М мЭіл©ФмЭЉлІМ мИШнЦЙ вАФ lesson_stageмЧР мУ∞мІА мХКмЭМ(мґ©лПМ л∞©мІА).
    upserts_disabled, reverts_disabled = upserts, reverts

    # 3) мҐЕнОЄ мШБмГБкЄЄмЭі
    if do_duration and durations:
        for lesson_id, (mt, path) in durations.items():
            sec = ffprobe_seconds(fs, path)
            if sec:
                sb.table("lessons").update({"duration_sec": sec}).eq("id", lesson_id).execute()
                updated_dur += 1

    # 4) мЮРлПЩ мЭіл©ФмЭЉ: мЭіл≤И мК§мЇФмЧРмДЬ мГИл°Ь мЩДл£М(waitвЖТdone)лРЬ лЛ®к≥Д вЖТ лЛ§мЭМ мВђмЪ©лЛ®к≥Д лЛілЛємЮРмЧРк≤М мХМл¶Љ
    emailed = notify_next(proj, enabled, lessons, cur, found, assignees, users, notified) if EMAIL_ENABLED else 0

    return {"ok": True, "project": proj["name"], "marked": 0, "found": len(upserts),
            "reverted": 0, "durations": updated_dur, "emailed": emailed}


def notify_next(proj, enabled, lessons, cur, found, assignees, users, notified):
    order = sorted(enabled, key=lambda s: stage_key(s["id"]))
    ids = [s["id"] for s in order]
    lmap = {l["id"]: l for l in lessons}
    sent = 0
    for (lesson_id, stage_id), (mtime, name) in found.items():
        prev = cur.get((lesson_id, stage_id))
        if prev and prev.get("status") == "done":
            continue                                   # мЭілѓЄ мЩДл£МмШАмЭМ вЖТ мХМл¶Љ мХИнХ®
        if (lesson_id, stage_id) in notified:
            continue                                   # мЭілѓЄ л∞ЬмЖ° кЄ∞л°Э мЮИмЭМ
        if stage_id not in ids:
            continue
        i = ids.index(stage_id)
        nxt = order[i + 1] if i + 1 < len(order) else None
        if not nxt:
            continue                                   # лІИмІАлІЙ лЛ®к≥Д вЖТ лЛ§мЭМ мЧЖмЭМ
        au = assignees.get(nxt["id"])
        user = users.get(au) if au else None
        to = (user or {}).get("email")
        lesson = lmap.get(lesson_id)
        lno = lesson.get("lesson_no") if lesson else "?"
        done_label = REF_LABEL.get(stage_id, str(stage_id))
        next_label = REF_LABEL.get(nxt["id"], str(nxt["id"]))
        subject = "[CDMS] %s ¬Ј %sм∞®мЛЬ вАФ '%s' м∞®л°АмЮЕлЛИлЛ§" % (proj["name"], lno, next_label)
        html = ("<div style='font-family:sans-serif;font-size:14px;color:#1c2430'>"
                "<p>мХИлЕХнХШмДЄмЪФ%s,</p>"
                "<p><b>%s</b> %sм∞®мЛЬмЭШ <b>%s</b> лЛ®к≥Дк∞А мЩДл£МлРШмЦі <b>%s</b> лЛ®к≥Дл•Љ мІДнЦЙнХ† м∞®л°АмЮЕлЛИлЛ§.</p>"
                "<p>м†ЬмЮС нШДнЩ©: <a href='%s'>%s</a></p>"
                "<p style='color:#8a94a6;font-size:12px'>вАФ CDMS мЮРлПЩмХМл¶Љ</p></div>") % (
                (" " + user["name"] if user and user.get("name") else ""),
                proj["name"], lno, done_label, next_label, CDMS_URL, CDMS_URL)
        if to:
            ok, err = send_email(to, subject, html)
            status = "sent" if ok else "error"
        else:
            ok, err, status = False, "лЛ§мЭМ лЛ®к≥Д(%s) лЛілЛємЮР лѓЄмІАм†Х" % next_label, "no_assignee"
        try:
            sb.table("email_notifications").upsert({
                "lesson_id": lesson_id, "completed_stage_id": stage_id, "next_stage_id": nxt["id"],
                "to_email": to, "to_user": au, "subject": subject,
                "status": status, "error": err}, on_conflict="lesson_id,completed_stage_id").execute()
        except Exception as e:
            log("email_notifications кЄ∞л°Э мЛ§нМ®:", e)
        if ok:
            sent += 1
        else:
            log("мЭіл©ФмЭЉ лѓЄл∞ЬмЖ°:", proj["name"], lno, next_label, "-", err)
    return sent


def action_probe_durations(fs, project_id):
    # кЄЄмЭі к≥ДмВ∞мЭА Edge(nas-versions scan)к∞А лЛілЛє вАФ мЧђкЄ∞мДЬлКФ мІДнЦЙ нСЬкЄ∞ мК§мЇФлІМ мИШнЦЙ
    return action_scan_progress(fs, project_id, do_duration=False)


# ============================================================================
# action: audio_check вАФ мҐЕнОЄ мШБмГБ мШ§лФФмШ§ нТИмІИ м†Рк≤А (лђімЭМ/нБіл¶ђнХС/к≥ЉлМА¬Јк≥ЉмЖМ мЭМлЯЙ)
#   к≤∞к≥ЉлКФ audio_checks нЕМмЭілЄФ(lesson_idлЛє 1нЦЙ upsert)мЧР кЄ∞л°Э вЖТ CDMS м∞®мЛЬ мГБмДЄмЧР нСЬмЛЬ
# ============================================================================
AUDIO_SILENCE_DB  = float(os.environ.get("AUDIO_SILENCE_DB", "-45"))   # лђімЭМ нМРм†Х мЮДк≥Д(dB)
AUDIO_SILENCE_MIN = float(os.environ.get("AUDIO_SILENCE_MIN", "2.0"))  # лђімЭМ мµЬмЖМ кЄЄмЭі(міИ)
AUDIO_LOUD_M      = float(os.environ.get("AUDIO_LOUD_M", "-9"))        # к≥ЉлМА мЭМлЯЙ(Momentary LUFS)
AUDIO_QUIET_M     = float(os.environ.get("AUDIO_QUIET_M", "-33"))      # к≥ЉмЖМ мЭМлЯЙ(Momentary LUFS)
AUDIO_JUMP_LU     = float(os.environ.get("AUDIO_JUMP_LU", "15"))       # кµђк∞Д к∞Д мЭМлЯЙ кЄЙл≥А мЮДк≥Д(LU) вАФ 9вЖТ15 мЩДнЩФ(2026-07-17, мШ§нГР к∞РмЖМ)
AUDIO_CH_MIN      = float(os.environ.get("AUDIO_CH_MIN", "3.0"))       # нХЬм™љ м±ДлДР лђімЭМ мµЬмЖМ кЄЄмЭі(міИ)
# мШБмГБ нТИмІИ кЄ∞м§А
VIDEO_W        = int(os.environ.get("VIDEO_W", "1920"))
VIDEO_H        = int(os.environ.get("VIDEO_H", "1080"))
VIDEO_KBPS_MIN = int(os.environ.get("VIDEO_KBPS_MIN", "3000"))
VIDEO_KBPS_MAX = int(os.environ.get("VIDEO_KBPS_MAX", "5500"))
VIDEO_FPS_MIN  = float(os.environ.get("VIDEO_FPS_MIN", "29.97"))
VIDEO_FPS_MAX  = float(os.environ.get("VIDEO_FPS_MAX", "30.0"))
VIDEO_BW_MIN   = float(os.environ.get("VIDEO_BW_MIN", "0.06"))         # лЄФлЮЩ/нЩФмЭінКЄ мµЬмЖМ мІАмЖН(міИ) вЙИ 2нФДл†ИмЮД@30fps
# л™ЕлПДлМАлєД(мЫєм†СкЈЉмД±) мЮРлПЩ к≤АмВђ
CCA_ENABLED  = os.environ.get("CCA_ENABLED", "true").lower() == "true"
CCA_INTERVAL = float(os.environ.get("CCA_INTERVAL", "5"))   # нФДл†ИмЮД мГШнФМ к∞Дк≤©(міИ)
CCA_RATIO    = float(os.environ.get("CCA_RATIO", "4.5"))    # кЄ∞м§А лМАлєД(WCAG AA)
CCA_MIN_H    = int(os.environ.get("CCA_MIN_H", "20"))       # к≤АмВђнХ† мµЬмЖМ кЄАмЮР лЖТмЭі(px) вАФ 14вЖТ20 мЩДнЩФ(2026-07-17)
CCA_MIN_CONF = int(os.environ.get("CCA_MIN_CONF", "82"))    # OCR мЛ†лҐ∞лПД нХШнХЬ вАФ 75вЖТ82 (мЮ•мЛЭ¬ЈмЫМнД∞лІИнБђ¬ЈмЭілѓЄмІА мШ§нГР к∞РмЖМ, 2026-07-18)
CCA_BG_STD   = float(os.environ.get("CCA_BG_STD", "12"))    # л∞∞к≤љ кЈ†мЭЉлПД(л∞ЭкЄ∞ нСЬм§АнОЄм∞®) мГБнХЬ вАФ мКђлЭЉмЭілУЬ¬ЈмЮРлІЙмЭА лЛ®мГЙ л∞∞к≤љ, мЭілѓЄмІА мЖН мЮ•мЛЭ нЕНмК§нКЄ л∞∞м†Ь

QC_VER = "v6-flatbg"  # нТИмІИ м†Рк≤А кЄ∞м§А л≤Дм†Д (v6: лЛ®мГЙ л∞∞к≤љ мЬД нХЩмКµ нЕНмК§нКЄлІМ к≤АмВђ вАФ мЭілѓЄмІА¬Јл™®мЕШкЈЄлЮШнФљ мЖН мЮ•мЛЭ нЕНмК§нКЄ м†ЬмЩЄ)


def _wcag_ratio(rgb1, rgb2):
    def lum(rgb):
        c = [v / 255.0 for v in rgb]
        c = [v / 12.92 if v <= 0.03928 else ((v + 0.055) / 1.055) ** 2.4 for v in c]
        return 0.2126 * c[0] + 0.7152 * c[1] + 0.0722 * c[2]
    l1, l2 = lum(rgb1), lum(rgb2)
    return (max(l1, l2) + 0.05) / (min(l1, l2) + 0.05)


def _analyze_contrast_local(local):
    """нФДл†ИмЮД мГШнФМ вЖТ OCR(tesseract)л°Ь нЕНмК§нКЄ мШБмЧ≠ к≤АмґЬ вЖТ кЄАмЮР/л∞∞к≤љмГЙ мґФм†Х вЖТ WCAG лМАлєД лѓЄлЛђ кµђк∞Д л∞ШнЩШ"""
    import shutil
    if not CCA_ENABLED:
        return []
    if not shutil.which("tesseract"):
        log("л™ЕлПДлМАлєД: tesseract лѓЄмД§мєШ вАФ к±ілДИлЬА (sudo apt-get install -y tesseract-ocr tesseract-ocr-kor)")
        return []
    try:
        from PIL import Image
    except ImportError:
        log("л™ЕлПДлМАлєД: pillow лѓЄмД§мєШ вАФ к±ілДИлЬА (pip install pillow)")
        return []
    langs = "kor+eng"
    try:
        o0 = subprocess.run(["tesseract", "--list-langs"], capture_output=True, text=True, timeout=30)
        if "kor" not in (o0.stdout or ""):
            langs = "eng"
    except Exception:
        langs = "eng"
    tmpd = tempfile.mkdtemp(prefix="cdms_cca_")
    issues = []
    try:
        subprocess.run(["ffmpeg", "-hide_banner", "-loglevel", "error", "-i", local, "-an",
                        "-vf", "fps=1/%g" % CCA_INTERVAL, os.path.join(tmpd, "f%05d.png")],
                       capture_output=True, timeout=3600)
        fails = []  # (t, (ratio, fg, bg, txt))
        for fn in sorted(os.listdir(tmpd)):
            try:
                idx = int(fn[1:6])
            except Exception:
                continue
            t = (idx - 1) * CCA_INTERVAL
            fp = os.path.join(tmpd, fn)
            o = subprocess.run(["tesseract", fp, "stdout", "-l", langs, "--psm", "11", "tsv"],
                               capture_output=True, text=True, timeout=120)
            img = None
            worst = None
            for ln in (o.stdout or "").splitlines()[1:]:
                cols = ln.split("\t")
                if len(cols) < 12:
                    continue
                try:
                    conf = float(cols[10]); txt = cols[11].strip()
                    x, y, w0, h0 = int(cols[6]), int(cols[7]), int(cols[8]), int(cols[9])
                except Exception:
                    continue
                if conf < CCA_MIN_CONF or len(txt) < 2 or h0 < CCA_MIN_H or w0 < CCA_MIN_H:
                    continue
                # нХЩмКµ нЕНмК§нКЄлІМ лМАмГБ: мЛ§м†Ь лђЄмЮР кµђмД± к≤АмВђ вАФ нХЬкЄА 2мЮРвЖС лШРлКФ мШБлђЄ¬ЈмИЂмЮР 3мЮРвЖС (мШЈ м£Љл¶Д¬ЈмЮ•мЛЭ мЪФмЖМмЭШ OCR нЧЫмЭЄмЛЭ л∞∞м†Ь)
                core = re.sub(r"[^0-9A-Za-zк∞А-нЮ£]", "", txt)
                hangul = re.sub(r"[^к∞А-нЮ£]", "", core)
                if not (len(hangul) >= 2 or len(core) >= 3):
                    continue
                if img is None:
                    img = Image.open(fp).convert("RGB")
                pad = max(2, h0 // 4)
                box = img.crop((max(0, x - pad), max(0, y - pad),
                                min(img.width, x + w0 + pad), min(img.height, y + h0 + pad)))
                px = list(box.getdata())
                if len(px) < 50:
                    continue
                lums = sorted((0.2126 * r + 0.7152 * g + 0.0722 * b, (r, g, b)) for r, g, b in px)
                n = len(lums)
                q = max(1, n // 4)
                q5 = max(1, n // 20)
                avg = lambda ps: tuple(sum(c[i] for c in ps) / len(ps) for i in range(3))
                # л∞∞к≤љ кЈ†мЭЉлПД: нХЩмКµ нЕНмК§нКЄ(мКђлЭЉмЭілУЬ¬ЈмЮРлІЙ)лКФ лЛ®мГЙ л∞∞к≤љ вАФ мЦілСРмЪі/л∞ЭмЭА м™љ м§С нХЬм™љмЭА кЈ†мЭЉнХімХЉ нХ®.
                #  мЦСм™љ л™®лСР мЦЉл£©лНЬл£©нХШл©і мЭілѓЄмІА¬Јл™®мЕШкЈЄлЮШнФљ мЖН лґАмИШм†Б нЕНмК§нКЄ(WCAG м†ЬмЩЄ лМАмГБ)л°Ь л≥ік≥† к±ілДИлЬА (2026-07-18)
                def _std(vs):
                    m = sum(vs) / len(vs)
                    return (sum((v - m) ** 2 for v in vs) / len(vs)) ** 0.5
                sd_dark = _std([l for l, _ in lums[:q]])
                sd_bright = _std([l for l, _ in lums[-q:]])
                if min(sd_dark, sd_bright) > CCA_BG_STD:
                    continue
                dark25, bright25 = avg([p for _, p in lums[:q]]), avg([p for _, p in lums[-q:]])
                dark5, bright5 = avg([p for _, p in lums[:q5]]), avg([p for _, p in lums[-q5:]])
                # мЛ§м†Ь кЄАмЮР нЪН мГЙ(кЈєлЛ® 5%) кЄ∞м§А вАФ к∞АлКФ кЄАмЮРмЧРмДЬ л∞∞к≤љ нФљмЕАмЭі мДЮмЧђ лМАлєДк∞А к≥ЉмЖМмЄ°м†ХлРШлНШ лђЄм†Ь мИШм†Х (2026-07-18)
                #  мЦілСРмЪі кЄАмЮР/л∞ЭмЭА л∞∞к≤љ, л∞ЭмЭА кЄАмЮР/мЦілСРмЪі л∞∞к≤љ мЦСл∞©нЦ• м§С нБ∞ к∞Т м±ДнГЭ
                ratio = max(_wcag_ratio(dark5, bright25), _wcag_ratio(dark25, bright5))
                # лМАлєД 1.25:1 лѓЄлІМмЭА кЄАмЮР/л∞∞к≤љмЭі мВђмЛ§мГБ к∞ЩмЭА мГЙ = OCRмЭі кЄАмЮРл•Љ мЭљмЧИмЭД мИШ мЧЖлКФ мШБмЧ≠(мЭілѓЄмІА мІИк∞Р мШ§мЭЄ) вЖТ м†ЬмЩЄ
                if ratio < 1.25:
                    continue
                if worst is None or ratio < worst[0]:
                    worst = (ratio, txt, (x, y, w0, h0), fp)
            if worst and worst[0] < CCA_RATIO:
                fails.append((t, worst))
        i = 0
        while i < len(fails):  # мЧ∞мЖН мЛ§нМ® нФДл†ИмЮД л≥СнХ©
            j = i
            while j + 1 < len(fails) and fails[j + 1][0] - fails[j][0] <= CCA_INTERVAL * 1.5:
                j += 1
            wr = min(fails[k][1][0] for k in range(i, j + 1))
            d = {"type": "contrast", "start": round(fails[i][0], 1),
                 "end": round(fails[j][0] + CCA_INTERVAL, 1),
                 "detail": "л™ЕлПДлМАлєД %.2f:1 вАФ кЄ∞м§А %.1f:1 лѓЄлЛђ (мЮРлІЙ/нЕНмК§нКЄ к∞АлПЕмД±)" % (wr, CCA_RATIO)}
            # лђЄм†Ь нФДл†ИмЮД мЇ°м≤Ш: кµђк∞Д лВі лМАлєДк∞А к∞АмЮ• лВЃмЭА нФДл†ИмЮДмЧР лђЄм†Ь мШБмЧ≠ л∞ХмК§ нСЬмЛЬ вЖТ previews/qc мЧЕл°ЬлУЬ
            try:
                import hashlib
                from PIL import ImageDraw
                k0 = min(range(i, j + 1), key=lambda k: fails[k][1][0])
                w0i = fails[k0][1]
                bb0 = w0i[2] if len(w0i) > 2 else None
                fp0 = w0i[3] if len(w0i) > 3 else None
                if fp0 and os.path.exists(fp0):
                    im0 = Image.open(fp0).convert("RGB")
                    if bb0:
                        dr = ImageDraw.Draw(im0)
                        x0, y0, w1, h1 = bb0
                        for off in range(3):
                            dr.rectangle([max(0, x0 - 6 - off), max(0, y0 - 6 - off),
                                          min(im0.width - 1, x0 + w1 + 6 + off), min(im0.height - 1, y0 + h1 + 6 + off)],
                                         outline=(255, 59, 48))
                    im0.thumbnail((960, 960))
                    ob = tempfile.NamedTemporaryFile(prefix="cdms_qc_", suffix=".jpg", delete=False)
                    obp = ob.name
                    ob.close()
                    im0.save(obp, "JPEG", quality=80)
                    with open(obp, "rb") as fj:
                        dat = fj.read()
                    try:
                        os.unlink(obp)
                    except Exception:
                        pass
                    key0 = hashlib.sha256(dat).hexdigest()[:24]
                    dest0 = "qc/%s.jpg" % key0
                    try:
                        sb.storage.from_("previews").upload(dest0, dat, {"content-type": "image/jpeg", "upsert": "true"})
                    except Exception as e2:
                        es2 = str(e2).lower()
                        if "exist" not in es2 and "duplicate" not in es2:
                            raise
                    d["shot"] = "%s/storage/v1/object/public/previews/%s" % (SB_URL, dest0)
            except Exception as e1:
                log("л™ЕлПДлМАлєД нФДл†ИмЮД мЇ°м≤Ш мЛ§нМ®:", e1)
            issues.append(d)
            i = j + 1
    except Exception as e:
        log("л™ЕлПДлМАлєД лґДмДЭ мЛ§нМ®:", e)
    finally:
        shutil.rmtree(tmpd, ignore_errors=True)
    return issues


def _analyze_video_local(local):
    """мШБмГБ нТИмІИ м†Рк≤А: кЈЬк≤©(нХімГБлПД/лєДнКЄл†ИмЭінКЄ/нФДл†ИмЮДл†ИмЭінКЄ) + лЄФлЮЩ/нЩФмЭінКЄ нФДл†ИмЮД кµђк∞Д"""
    issues, stats = [], {}
    # 1) кЈЬк≤© (ffprobe)
    po = subprocess.run(["ffprobe", "-v", "error", "-select_streams", "v:0",
                         "-show_entries", "stream=width,height,avg_frame_rate,bit_rate",
                         "-show_entries", "format=bit_rate", "-of", "json", local],
                        capture_output=True, text=True, timeout=120)
    j = json.loads(po.stdout or "{}")
    st0 = (j.get("streams") or [{}])[0]
    w, h = st0.get("width"), st0.get("height")
    afr = st0.get("avg_frame_rate") or "0/1"
    try:
        num, den = afr.split("/")
        fps = (float(num) / float(den)) if float(den) else 0.0
    except Exception:
        fps = 0.0
    br = st0.get("bit_rate") or (j.get("format") or {}).get("bit_rate") or 0
    try:
        kbps = int(int(br) / 1000)
    except Exception:
        kbps = 0
    stats["video"] = {"w": w, "h": h, "fps": round(fps, 2), "kbps": kbps}
    if w and h and (int(w), int(h)) != (VIDEO_W, VIDEO_H):
        issues.append({"type": "spec", "start": 0, "end": 0,
                       "detail": "нФДл†ИмЮД мВђмЭім¶И {}√Ч{} вАФ кЄ∞м§А {}√Ч{}".format(w, h, VIDEO_W, VIDEO_H)})
    if kbps and not (VIDEO_KBPS_MIN <= kbps <= VIDEO_KBPS_MAX):
        issues.append({"type": "spec", "start": 0, "end": 0,
                       "detail": "лєДнКЄм†ДмЖ°л•† {:,}kbps вАФ кЄ∞м§А {:,}~{:,}kbps".format(kbps, VIDEO_KBPS_MIN, VIDEO_KBPS_MAX)})
    if fps and not (VIDEO_FPS_MIN - 0.01 <= fps <= VIDEO_FPS_MAX + 0.01):
        issues.append({"type": "spec", "start": 0, "end": 0,
                       "detail": "нФДл†ИмЮДл†ИмЭінКЄ %.2ffps вАФ кЄ∞м§А %g~%gfps" % (fps, VIDEO_FPS_MIN, VIDEO_FPS_MAX)})
    # 2) 100% лЄФлЮЩ/нЩФмЭінКЄ нФДл†ИмЮД (2нФДл†ИмЮД мЭімГБ мЧ∞мЖН)
    bd = "blackdetect=d=%g:pic_th=0.98:pix_th=0.05" % VIDEO_BW_MIN
    def detect(filters):
        out = subprocess.run(["ffmpeg", "-hide_banner", "-nostats", "-i", local, "-an",
                              "-vf", filters, "-f", "null", "-"],
                             capture_output=True, text=True, timeout=3600)
        segs = []
        for mm in re.finditer(r"black_start:\s*(-?\d+\.?\d*)\s+black_end:\s*(-?\d+\.?\d*)", out.stderr or ""):
            segs.append((float(mm.group(1)), float(mm.group(2))))
        return segs
    for a, b2 in detect(bd):
        issues.append({"type": "black", "start": round(a, 2), "end": round(b2, 2),
                       "detail": "%.2fміИ лЄФлЮЩ(к≤АмЭА нЩФл©і) нФДл†ИмЮД" % (b2 - a)})
    for a, b2 in detect("negate," + bd):
        issues.append({"type": "white", "start": round(a, 2), "end": round(b2, 2),
                       "detail": "%.2fміИ нЩФмЭінКЄ(нЭ∞ нЩФл©і) нФДл†ИмЮД" % (b2 - a)})
    return issues, stats


def _analyze_audio(fs, remote_path):
    """ffmpeg 1нЪМ мЛ§нЦЙ(silencedetect+ebur128+volumedetect) вЖТ (issues, stats)"""
    local, is_tmp = fs.local_copy(remote_path)
    try:
        af = "silencedetect=noise=%gdB:d=%g,ebur128=metadata=0,volumedetect" % (AUDIO_SILENCE_DB, AUDIO_SILENCE_MIN)
        out = subprocess.run(
            ["ffmpeg", "-hide_banner", "-nostats", "-i", local, "-vn", "-map", "0:a:0?",
             "-af", af, "-f", "null", "-"],
            capture_output=True, text=True, timeout=1800)
        err = out.stderr or ""
        issues, stats = [], {}
        m = re.search(r"Duration:\s*(\d+):(\d+):(\d+\.?\d*)", err)
        dur = int(m.group(1)) * 3600 + int(m.group(2)) * 60 + float(m.group(3)) if m else None
        stats["duration"] = dur
        if "does not contain any stream" in err or "matches no streams" in err:
            issues.append({"type": "no_audio", "start": 0, "end": dur or 0, "detail": "мШ§лФФмШ§ нКЄлЮЩ мЧЖмЭМ"})
            return issues, stats
        # лђімЭМ кµђк∞Д
        starts = [float(x) for x in re.findall(r"silence_start:\s*(-?\d+\.?\d*)", err)]
        ends = [(float(a), float(b)) for a, b in
                re.findall(r"silence_end:\s*(-?\d+\.?\d*)\s*\|\s*silence_duration:\s*(-?\d+\.?\d*)", err)]
        for i, st in enumerate(starts):
            if i < len(ends):
                en, d = ends[i]
            else:  # нММмЭЉ лБЭкєМмІА лђімЭМ
                en = dur or st
                d = en - st
            issues.append({"type": "silence", "start": round(max(0, st), 1), "end": round(en, 1),
                           "detail": "%.1fміИ лђімЭМ" % d})
        # м†Дм≤і л≥Љл•® / нБіл¶ђнХС
        m = re.search(r"mean_volume:\s*(-?\d+\.?\d*)\s*dB", err)
        if m: stats["mean"] = float(m.group(1))
        m = re.search(r"max_volume:\s*(-?\d+\.?\d*)\s*dB", err)
        if m: stats["max"] = float(m.group(1))
        if stats.get("max") is not None and stats["max"] >= -0.1:
            issues.append({"type": "clip", "start": 0, "end": round(dur or 0, 1),
                           "detail": "мµЬлМА л≥Љл•® %.1fdB вАФ нБіл¶ђнХС(мЭМ кє®мІР) мЭШмЛђ" % stats["max"]})
        # кµђк∞Д лЭЉмЪ∞лУЬлЛИмК§(ebur128 Momentary) вЖТ к≥ЉлМА/к≥ЉмЖМ
        frames3 = [(float(t), float(m2), float(s2)) for t, m2, s2 in
                   re.findall(r"t:\s*(-?\d+\.?\d*)\s+TARGET:.*?M:\s*(-?\d+\.?\d*)\s+S:\s*(-?\d+\.?\d*)", err)]
        frames = [(t, m2) for t, m2, s2 in frames3]
        def runs(flag_fn, min_len):
            segs, s0, last = [], None, None
            for t, mv in frames:
                if flag_fn(mv):
                    if s0 is None: s0 = t
                    last = t
                else:
                    if s0 is not None and last is not None and last - s0 >= min_len:
                        segs.append((s0, last))
                    s0 = None
            if s0 is not None and last is not None and last - s0 >= min_len:
                segs.append((s0, last))
            return segs
        sil = [(i["start"], i["end"]) for i in issues if i["type"] == "silence"]
        def in_sil(a, b):
            return any(not (b < s or a > e) for s, e in sil)
        for a, b in runs(lambda v: v >= AUDIO_LOUD_M, 1.0):
            issues.append({"type": "loud", "start": round(a, 1), "end": round(b, 1),
                           "detail": "мЭМлЯЙ к≥ЉлМА вАФ мИЬк∞Д лЭЉмЪ∞лУЬлЛИмК§ %g LUFS мЭімГБ" % AUDIO_LOUD_M})
        for a, b in runs(lambda v: -70 < v <= AUDIO_QUIET_M, 3.0):
            if not in_sil(a, b):
                issues.append({"type": "quiet", "start": round(a, 1), "end": round(b, 1),
                               "detail": "мЭМлЯЙ к≥ЉмЖМ вАФ лМАмВђк∞А мЮШ мХИ лУ§л¶і мИШ мЮИмЭМ"})
        # кµђк∞Д к∞Д мЭМлЯЙ кЄЙл≥А (Short-term лЭЉмЪ∞лУЬлЛИмК§к∞А 3міИ м†Д лМАлєД AUDIO_JUMP_LU мЭімГБ м∞®мЭі)
        svals = [(t, s2) for t, m2, s2 in frames3 if s2 > -70]
        marks = []
        for i in range(30, len(svals)):
            t1, s1 = svals[i]
            t0, s0 = svals[i - 30]
            if t1 - t0 <= 4.0 and abs(s1 - s0) >= AUDIO_JUMP_LU:
                marks.append((t1, s1 - s0))
        i = 0
        while i < len(marks):
            j = i
            while j + 1 < len(marks) and marks[j + 1][0] - marks[j][0] <= 2.0:
                j += 1
            a, b = marks[i][0], marks[j][0]
            mx = max(abs(d) for _, d in marks[i:j + 1])
            if not in_sil(max(0, a - 3), b):
                issues.append({"type": "jump", "start": round(max(0, a - 3), 1), "end": round(b, 1),
                               "detail": "мЭМлЯЙ кЄЙл≥А вАФ кµђк∞Д к∞Д %.0f LU м∞®мЭі" % mx})
            i = j + 1
        # мК§нЕМл†ИмШ§ нХЬм™љ м±ДлДР лђімЭМ (м±ДлДРл≥Д silencedetect нЫД нХЬм™љлІМ лђімЭМмЭЄ кµђк∞Д)
        ch = 0
        try:
            po = subprocess.run(["ffprobe", "-v", "error", "-select_streams", "a:0",
                                 "-show_entries", "stream=channels", "-of", "default=nw=1:nokey=1", local],
                                capture_output=True, text=True, timeout=60)
            ch = int((po.stdout or "0").strip() or 0)
        except Exception:
            ch = 0
        if ch >= 2:
            def ch_sil(cidx):
                o = subprocess.run(["ffmpeg", "-hide_banner", "-nostats", "-i", local, "-vn", "-map", "0:a:0",
                                    "-af", "pan=mono|c0=c%d,silencedetect=noise=-50dB:d=%g" % (cidx, AUDIO_CH_MIN),
                                    "-f", "null", "-"], capture_output=True, text=True, timeout=1800)
                e2 = o.stderr or ""
                ss = [float(x) for x in re.findall(r"silence_start:\s*(-?\d+\.?\d*)", e2)]
                ee = [float(a) for a, b in
                      re.findall(r"silence_end:\s*(-?\d+\.?\d*)\s*\|\s*silence_duration:\s*(-?\d+\.?\d*)", e2)]
                segs = []
                for k, s0 in enumerate(ss):
                    e0 = ee[k] if k < len(ee) else (dur or s0)
                    segs.append((max(0, s0), e0))
                return segs

            def sub_ivals(a_list, b_list):
                out = []
                for a0, a1 in a_list:
                    cur = [(a0, a1)]
                    for b0, b1 in b_list:
                        nxt = []
                        for c0, c1 in cur:
                            if b1 <= c0 or b0 >= c1:
                                nxt.append((c0, c1)); continue
                            if b0 > c0: nxt.append((c0, b0))
                            if b1 < c1: nxt.append((b1, c1))
                        cur = nxt
                    out.extend(cur)
                return [(x, y) for x, y in out if y - x >= AUDIO_CH_MIN]

            L, R = ch_sil(0), ch_sil(1)
            for a, b in sub_ivals(L, R):
                issues.append({"type": "channel", "start": round(a, 1), "end": round(b, 1),
                               "detail": "мЩЉм™љ(L) м±ДлДР лђімЭМ вАФ мШ§л•Єм™љлІМ мґЬл†•"})
            for a, b in sub_ivals(R, L):
                issues.append({"type": "channel", "start": round(a, 1), "end": round(b, 1),
                               "detail": "мШ§л•Єм™љ(R) м±ДлДР лђімЭМ вАФ мЩЉм™љлІМ мґЬл†•"})
        # вФАвФА мШБмГБ нТИмІИ м†Рк≤А (кЈЬк≤© + лЄФлЮЩ/нЩФмЭінКЄ нФДл†ИмЮД) вФАвФА
        try:
            vissues, vstats = _analyze_video_local(local)
            issues.extend(vissues)
            stats.update(vstats)
        except Exception as e:
            log("мШБмГБ нТИмІИ м†Рк≤А мЛ§нМ®:", e)
        try:
            issues.extend(_analyze_contrast_local(local))
        except Exception as e:
            log("л™ЕлПДлМАлєД м†Рк≤А мЛ§нМ®:", e)
        issues.sort(key=lambda x: x["start"])
        return issues, stats
    finally:
        if is_tmp:
            try: os.unlink(local)
            except Exception: pass


def _qc_basekey(nm):
    """нТИмІИ м†Рк≤АмЪ© нММмЭЉ кЄ∞л≥ЄнВ§ вАФ л¶ђлєДм†Д¬Ј(мИШм†Х)¬Јл≤Дм†Д нСЬкЄ∞л•Љ м†Ьк±∞нХі к∞ЩмЭА нММнКЄмЭШ мВђл≥Є/мИШм†Хл≥ЄмЭД лђґлКФлЛ§"""
    s = re.sub(r"\.[A-Za-z0-9]+$", "", nm or "")
    s = re.sub(r"(?i)re\s*\d+|v\d+(\.\d+)*|\(\d+\)|\(мИШм†Х\)|мИШм†Хл≥Є?|мµЬмҐЕ", "", s)
    return re.sub(r"\s+", "", s).lower()


def action_audio_check(fs, project_id, lesson_id=None, notify_user=None):
    b = load_project_bundle(project_id)
    if not b:
        return {"ok": False, "error": "project not found"}
    proj, enabled, lessons = b["proj"], b["enabled"], b["lessons"]
    root = proj.get("nas_root")
    if not root:
        return {"ok": False, "error": "nas_root лѓЄмД§м†Х вАФ л®Љм†А NAS нПілНФ мІАм†ХмЭі нХДмЪФнХ©лЛИлЛ§"}
    mroot = re.match(r"^nas(\d+):(.*)$", root)
    if mroot:  # лЛ§л•Є NAS: м∞®мЛЬ лЛ®мЬДлІМ мЫРк≤© м†Рк≤А(мШБмГБ лЛ§мЪіл°ЬлУЬ нЫД лґДмДЭ)
        if not lesson_id:
            return {"ok": False, "error": "лЛ§л•Є NAS(nasN:) к≥Љм†ХмЭА м∞®мЛЬ лЛ®мЬД м†Рк≤АлІМ мІАмЫРнХ©лЛИлЛ§ вАФ м∞®мЛЬ мГБмДЄмЭШ 'рЯФК мЭі м∞®мЛЬ нТИмІИ м†Рк≤А' л≤ДнКЉмЭД мВђмЪ©нХШмДЄмЪФ"}
        return _audio_check_remote(int(mroot.group(1)), mroot.group(2), b, lesson_id, notify_user)
    has_weeks = any(l.get("week_no") for l in lessons)
    st = next((s for s in enabled if s["id"] == LENGTH_STAGE_ID), None)
    folder = "%s/%s" % (root, st["nas_folder"]) if st else None
    targets = {}  # lesson_id -> {basekey: (mtime, path, name)} вАФ м∞®мЛЬ лВі "л™®лУ†" мҐЕнОЄ нММнКЄ(мВђл≥Є¬ЈмИШм†Хл≥ЄмЭА мµЬмЛ† 1к∞Ь)
    if folder and fs.exists(folder):
        for rel, mtime in fs.walkfiles(folder):
            base = rel.rsplit("/", 1)[-1]
            if ".cdms_" in rel:
                continue
            if base.startswith("~") or os.path.splitext(base)[1].lower() not in VIDEO_EXT:
                continue
            l = match_lesson(rel, lessons, has_weeks)
            if not l:
                continue
            if lesson_id and l["id"] != lesson_id:
                continue
            grp = targets.setdefault(l["id"], {})
            k = _qc_basekey(base)
            if k not in grp or mtime >= grp[k][0]:
                grp[k] = (mtime, "%s/%s" % (folder, rel), base)
    if not targets:
        return {"ok": False, "error": "мҐЕнОЄ нПілНФмЧРмДЬ м∞®мЛЬмЧР лІ§мє≠лРШлКФ мШБмГБмЭД м∞ЊмІА л™їнЦИмКµлЛИлЛ§"}
    checked, problems = 0, 0
    now = datetime.now(timezone.utc).isoformat()
    for lid, grp in targets.items():
        names = []
        exist = {}
        try:  # м¶ЭлґД м†Рк≤А: нММмЭЉмЭі мХИ л∞ФлАМк≥† кЄ∞м§АлПД к∞ЩмЬЉл©і мЮђлґДмДЭ мГЭлЮµ
            exist = {r0["file_name"]: r0 for r0 in (sb.table("audio_checks").select("file_name,file_mtime,qc_ver,status").eq("lesson_id", lid).execute().data or [])}
        except Exception:
            pass
        for (mt, path, name) in sorted(grp.values(), key=lambda x: x[2]):
            e0 = exist.get(name)
            if e0 and e0.get("file_mtime") == int(mt or 0) and e0.get("qc_ver") == QC_VER and e0.get("status") in ("ok", "warn", "bad"):
                names.append(name)
                checked += 1
                log("нТИмІИ м†Рк≤А(мГЭлЮµвАФл≥Ак≤љ мЧЖмЭМ):", name)
                continue
            row = {"lesson_id": lid, "project_id": project_id, "file_name": name, "file_path": path,
                   "file_mtime": int(mt or 0), "qc_ver": QC_VER, "checked_at": now}
            try:
                issues, stats = _analyze_audio(fs, path)
                bad = any(i["type"] in ("silence", "clip", "no_audio", "channel") for i in issues)
                warn = any(i["type"] in ("loud", "quiet", "jump", "spec", "black", "white", "contrast") for i in issues)
                row.update({"duration_sec": stats.get("duration"), "mean_volume": stats.get("mean"),
                            "max_volume": stats.get("max"), "issues": issues,
                            "status": "bad" if bad else ("warn" if warn else "ok"), "error": None})
                problems += len(issues)
                checked += 1
            except Exception as e:
                row.update({"status": "error", "error": str(e), "issues": []})
            sb.table("audio_checks").upsert(row, on_conflict="lesson_id,file_name").execute()
            names.append(name)
            log("нТИмІИ м†Рк≤А:", name, row.get("status"))
        # мЭіл≤И м†Рк≤А лМАмГБмЧР мЧЖлКФ нММмЭЉмЭШ мШИм†Д к≤∞к≥Љ м†Хл¶ђ(нММмЭЉ кµРм≤і¬ЈмЭіл¶Д л≥Ак≤љ лМАмЭС)
        try:
            old = sb.table("audio_checks").select("id,file_name").eq("lesson_id", lid).execute().data or []
            for o in old:
                if o.get("file_name") not in names:
                    sb.table("audio_checks").delete().eq("id", o["id"]).execute()
        except Exception:
            pass

    # мЧЕл°ЬлУЬ лЛілЛємЮР мЭіл©ФмЭЉ мХМл¶Љ (лђЄм†Ь л∞Ьк≤ђ мЛЬ)
    emailed = 0
    if notify_user and problems > 0 and EMAIL_ENABLED:
        try:
            u = sb.table("users").select("name,email").eq("id", notify_user).execute().data
            email = u and u[0].get("email")
            if email:
                TYPE_KR = {"silence": "лђімЭМ", "clip": "нБіл¶ђнХС", "loud": "к≥ЉлМАмЭМлЯЙ", "quiet": "к≥ЉмЖМмЭМлЯЙ",
                           "no_audio": "мШ§лФФмШ§ мЧЖмЭМ", "jump": "мЭМлЯЙ кЄЙл≥А", "channel": "м±ДлДР лђімЭМ"}
                lmap = {l["id"]: l for l in lessons}
                def mmss(v):
                    v = int(float(v or 0)); return "%d:%02d" % (v // 60, v % 60)
                res = sb.table("audio_checks").select("*").in_("lesson_id", list(targets.keys())).execute().data
                blocks = []
                for r0 in (res or []):
                    iss = r0.get("issues") or []
                    if not iss:
                        continue
                    l0 = lmap.get(r0["lesson_id"]) or {}
                    items = "".join(
                        "<li>%s ~ %s вАФ <b>%s</b> ¬Ј %s</li>" %
                        (mmss(i.get("start")), mmss(i.get("end")),
                         TYPE_KR.get(i.get("type"), i.get("type")), i.get("detail") or "")
                        for i in iss)
                    blocks.append("<p style='margin:12px 0 4px'><b>%sм∞®мЛЬ</b> %s ¬Ј нММмЭЉ: %s</p><ul style='margin:4px 0'>%s</ul>"
                                  % (l0.get("lesson_no", ""), l0.get("title") or "", r0.get("file_name") or "", items))
                if blocks:
                    html = ("<div style=\"font-family:Apple SD Gothic Neo,Malgun Gothic,sans-serif;font-size:14px;color:#222;line-height:1.6\">"
                            "<p>мХИлЕХнХШмДЄмЪФ, CDMS мШ§лФФмШ§ м†Рк≤А мХМл¶ЉмЮЕлЛИлЛ§.</p>"
                            "<p>мЧЕл°ЬлУЬнХШмЛ† <b>%s</b> мҐЕнОЄ мШБмГБмЧРмДЬ <b>нТИмІИ лђЄм†Ь %dк±і</b>мЭі л∞Ьк≤ђлРШмЧИмКµлЛИлЛ§.</p>" % (proj["name"], problems)
                            + "".join(blocks) +
                            "<p><a href=\"%s\" style=\"display:inline-block;background:#4b3fbb;color:#fff;text-decoration:none;padding:8px 16px;border-radius:8px\">CDMSмЧРмДЬ нЩХмЭЄнХШкЄ∞</a></p>" % CDMS_URL +
                            "<p style=\"color:#999;font-size:12px\">CDMSмЧРмДЬ нХілЛє м∞®мЛЬл•Љ нБіл¶≠нХШл©і 'рЯФК нТИмІИ м†Рк≤А' мДємЕШмЧР лђЄм†Ь кµђк∞ДмЭі нСЬмЛЬлРШл©∞, мЛЬк∞ДмЭД нБіл¶≠нХШл©і кЈЄ мЬДмєШлґАнД∞ мЮђмГЭлР©лЛИлЛ§.</p></div>")
                    send_email(email, "[CDMS] нТИмІИ м†Рк≤А к≤∞к≥Љ вАФ %s лђЄм†Ь %dк±і" % (proj["name"], problems), html)
                    emailed = 1
        except Exception as e:
            log("мШ§лФФмШ§ мХМл¶Љ л©ФмЭЉ мЛ§нМ®:", e)
    return {"ok": True, "checked": checked, "problems": problems, "emailed": emailed}


# ============================================================================
# action: scan_file вАФ CDMS мЧЕл°ЬлУЬ нММмЭЉ л∞±мЛ† к≤АмВђ(ClamAV) нЫД мµЬмҐЕ нПілНФл°Ь мЭілПЩ
#   мЧЕл°ЬлУЬлКФ .cdms_scan(к≤АмВђ лМАкЄ∞)мЧР м†АмЮ•лР® вЖТ нЖµк≥Љ: destл°Ь мЭілПЩ / к∞РмЧЉ: .cdms_blocked к≤©л¶ђ + л©ФмЭЉ
#   ClamAV лѓЄмД§мєШ мЛЬ мЛ§нМ®(fail-closed): нММмЭЉмЭА мµЬмҐЕ нПілНФмЧР л∞ШмШБлРШмІА мХКмЭМ
# ============================================================================
def _clam_run(localpath):
    """ClamAV мЛ§нЦЙ вЖТ (rc, нГРмІАл™Е, мШ§л•ШлђЄмЮРмЧі). лѓЄмД§мєШ мЛЬ rc=None"""
    import shutil
    scanner = shutil.which("clamdscan") or shutil.which("clamscan")
    if not scanner:
        return None, "", "мДЬл≤ДмЧР ClamAVк∞А мЧЖмКµлЛИлЛ§. sudo apt-get install -y clamav clamav-daemon нЫД мЮђмЛЬлПД"
    args = [scanner, "--no-summary"]
    if scanner.endswith("clamdscan"):
        args.append("--fdpass")
    out = subprocess.run(args + [localpath], capture_output=True, text=True, timeout=3600)
    sig = ""
    for ln in (out.stdout or "").splitlines():
        if ": " in ln and not ln.rstrip().endswith("OK"):
            sig = ln.split(": ", 1)[1].strip()
            break
    return out.returncode, sig, (out.stderr or "")[:200]


def _blocked_email(name, sig, qpath, notify):
    if not EMAIL_ENABLED:
        return
    try:
        emails = []
        if notify:
            u = sb.table("users").select("email").eq("id", notify).execute().data
            if u and u[0].get("email"): emails.append(u[0]["email"])
        ad = sb.table("user_roles").select("user_id,users:user_id(email)").eq("role_code", "admin").execute().data
        for a in (ad or []):
            e = (a.get("users") or {}).get("email")
            if e: emails.append(e)
        html = ("<div style='font-family:Malgun Gothic,sans-serif;font-size:14px;line-height:1.6'>"
                "<p>CDMSл°Ь мЧЕл°ЬлУЬлРЬ нММмЭЉмЧРмДЬ <b>мХЕмД±мљФлУЬк∞А нГРмІАлРШмЦі м∞®лЛ®</b>лРШмЧИмКµлЛИлЛ§.</p>"
                "<p>нММмЭЉ: <b>%s</b><br>нГРмІАл™Е: %s<br>к≤©л¶ђ мЬДмєШ: %s</p>"
                "<p>мЭі нММмЭЉмЭА NAS мЮСмЧЕ нПілНФмЧР л∞ШмШБлРШмІА мХКмХШмКµлЛИлЛ§. мЧЕл°ЬлУЬнХЬ PCмЭШ л∞±мЛ† м†Рк≤АмЭД кґМмЮ•нХ©лЛИлЛ§.</p></div>"
                % (name, sig or "-", qpath))
        for em in sorted(set(emails)):
            send_email(em, "[CDMS] вЫФ мЧЕл°ЬлУЬ м∞®лЛ® вАФ мХЕмД±мљФлУЬ нГРмІА: %s" % name, html)
    except Exception as e:
        log("м∞®лЛ® л©ФмЭЉ мЛ§нМ®:", e)


# вФАвФА мЫРк≤© NAS(nasN:) FileStation API нЧђнНЉ вАФ NAS2 лУ± лІИмЪінКЄ мХИ лРЬ NASмЭШ к≤АмВђмЪ© вФАвФА
def _syno_api(base, path_, params, timeout=120):
    import urllib.request, urllib.parse
    url = base + path_ + "?" + urllib.parse.urlencode(params)
    with urllib.request.urlopen(url, timeout=timeout) as r:
        return json.loads(r.read().decode("utf-8", "ignore"))


def _syno_login2(cfg):
    base = (cfg.get("url") or "").rstrip("/")
    j = _syno_api(base, "/webapi/auth.cgi", {"api": "SYNO.API.Auth", "version": "6", "method": "login",
                                             "account": cfg.get("username") or "", "passwd": cfg.get("password") or "",
                                             "session": "FileStation", "format": "sid"})
    if not j.get("success"):
        raise RuntimeError("NAS л°ЬкЈЄмЭЄ мЛ§нМ® code=%s" % ((j.get("error") or {}).get("code")))
    return base, j["data"]["sid"]


def _syno_ls(base, sid, folder):
    j = _syno_api(base, "/webapi/entry.cgi", {"api": "SYNO.FileStation.List", "version": "2", "method": "list",
                                              "folder_path": json.dumps(folder), "_sid": sid})
    return ((j.get("data") or {}).get("files")) or []


def _syno_download(base, sid, path, dst):
    import urllib.request, urllib.parse
    url = base + "/webapi/entry.cgi?" + urllib.parse.urlencode(
        {"api": "SYNO.FileStation.Download", "version": "2", "method": "download", "mode": "open",
         "path": path, "_sid": sid})
    with urllib.request.urlopen(url, timeout=3600) as r, open(dst, "wb") as f:
        while True:
            b = r.read(1024 * 1024)
            if not b: break
            f.write(b)


def _syno_move(base, sid, src, dest_folder):
    j = _syno_api(base, "/webapi/entry.cgi", {"api": "SYNO.FileStation.CopyMove", "version": "3", "method": "start",
                                              "path": json.dumps([src]), "dest_folder_path": json.dumps(dest_folder),
                                              "remove_src": "true", "_sid": sid})
    if not j.get("success"):
        raise RuntimeError("мЭілПЩ мЛЬмЮС мЛ§нМ® code=%s" % ((j.get("error") or {}).get("code")))
    tid = j["data"]["taskid"]
    for _ in range(180):
        s = _syno_api(base, "/webapi/entry.cgi", {"api": "SYNO.FileStation.CopyMove", "version": "3",
                                                  "method": "status", "taskid": json.dumps(tid), "_sid": sid})
        if (s.get("data") or {}).get("finished"):
            return True
        time.sleep(1)
    raise RuntimeError("мЭілПЩ мЛЬк∞Д міИк≥Љ")


def _syno_rename(base, sid, path, newname):
    j = _syno_api(base, "/webapi/entry.cgi", {"api": "SYNO.FileStation.Rename", "version": "2", "method": "rename",
                                              "path": json.dumps([path]), "name": json.dumps([newname]), "_sid": sid})
    if not j.get("success"):
        raise RuntimeError("мЭіл¶Дл≥Ак≤љ мЛ§нМ® code=%s" % ((j.get("error") or {}).get("code")))


def _syno_mkdir(base, sid, parent, name):
    try:
        _syno_api(base, "/webapi/entry.cgi", {"api": "SYNO.FileStation.CreateFolder", "version": "2", "method": "create",
                                              "folder_path": json.dumps([parent]), "name": json.dumps([name]),
                                              "force_parent": "true", "_sid": sid})
    except Exception:
        pass


def _scan_file_remote(nid, path, dest, name, notify):
    rows = sb.table("nas_config").select("*").eq("id", nid).execute().data
    if not rows:
        return {"ok": False, "error": "nas_config(id=%s) мЧЖмЭМ" % nid}
    base, sid = _syno_login2(rows[0])
    try:
        folder = path.rsplit("/", 1)[0]
        found = False
        for _ in range(20):
            if any(f.get("name") == name for f in _syno_ls(base, sid, folder)):
                found = True
                break
            time.sleep(3)
        if not found:
            return {"ok": False, "error": "к≤АмВђ лМАмГБ нММмЭЉмЭД м∞ЊмІА л™їнЦИмКµлЛИлЛ§: nas%s:%s" % (nid, path)}
        tf = tempfile.NamedTemporaryFile(prefix="cdms_scan_", delete=False)
        tmp = tf.name
        tf.close()
        try:
            _syno_download(base, sid, path, tmp)
            rc, sig, err = _clam_run(tmp)
        finally:
            try: os.unlink(tmp)
            except Exception: pass
        if rc is None:
            return {"ok": False, "error": err}
        if rc == 0:
            names = {f.get("name") for f in _syno_ls(base, sid, dest)}
            b0, e0 = os.path.splitext(name)
            fname, n = name, 1
            while fname in names:
                n += 1
                fname = "%s(%d)%s" % (b0, n, e0)
            path2 = path
            if fname != name:
                _syno_rename(base, sid, path, fname)
                path2 = folder + "/" + fname
            _syno_move(base, sid, path2, dest)
            log("л∞±мЛ†к≤АмВђ нЖµк≥Љ(nas%s):" % nid, dest + "/" + fname)
            return {"ok": True, "clean": True, "moved": "nas%s:%s/%s" % (nid, dest, fname)}
        if rc == 1:
            stage = folder.rsplit("/", 1)[0]
            _syno_mkdir(base, sid, stage, ".cdms_blocked")
            qpath = "nas%s:%s/.cdms_blocked/%s" % (nid, stage, name)
            try: _syno_move(base, sid, path, stage + "/.cdms_blocked")
            except Exception: qpath = "nas%s:%s" % (nid, path)
            log("мХЕмД±мљФлУЬ м∞®лЛ®(nas%s):" % nid, name, sig)
            _blocked_email(name, sig, qpath, notify)
            return {"ok": True, "clean": False, "virus": sig or "malware", "quarantine": qpath}
        return {"ok": False, "error": "к≤АмВђ мШ§л•Ш(rc=%s): %s" % (rc, err)}
    finally:
        try:
            _syno_api(base, "/webapi/auth.cgi", {"api": "SYNO.API.Auth", "version": "6", "method": "logout",
                                                 "session": "FileStation", "_sid": sid})
        except Exception:
            pass


class _RemoteFSLite:
    """мЫРк≤© NAS нММмЭЉмЭД мЮДмЛЬ лЛ§мЪіл°ЬлУЬнХі fs.local_copyм≤ШлЯЉ м†Ьк≥µ (мШ§лФФмШ§ лґДмДЭмЪ©)"""
    def __init__(self, base, sid):
        self.base, self.sid = base, sid
    def local_copy(self, path):
        tf = tempfile.NamedTemporaryFile(prefix="cdms_aud_", delete=False)
        tmp = tf.name
        tf.close()
        _syno_download(self.base, self.sid, path, tmp)
        return tmp, True


def _syno_walk(base, sid, folder, depth):
    out = []
    j = _syno_api(base, "/webapi/entry.cgi", {"api": "SYNO.FileStation.List", "version": "2", "method": "list",
                                              "folder_path": json.dumps(folder), "additional": json.dumps(["time"]),
                                              "_sid": sid})
    for f in ((j.get("data") or {}).get("files")) or []:
        nm = f.get("name") or ""
        if nm.startswith(".cdms_") or "#recycle" in nm.lower() or nm.lower() == "old":
            continue
        if f.get("isdir"):
            if depth > 0:
                out.extend(_syno_walk(base, sid, f["path"], depth - 1))
        else:
            mt = ((f.get("additional") or {}).get("time") or {}).get("mtime") or 0
            out.append((nm, f["path"], mt))
    return out


def _audio_check_remote(nid, root, bundle, lesson_id, notify):
    """лЛ§л•Є NAS(nasN:) к≥Љм†ХмЭШ м∞®мЛЬ лЛ®мЬД мШ§лФФмШ§ м†Рк≤А вАФ мШБмГБмЭД лВіл†§л∞ЫмХД лґДмДЭ"""
    rows = sb.table("nas_config").select("*").eq("id", nid).execute().data
    if not rows:
        return {"ok": False, "error": "nas_config(id=%s) мЧЖмЭМ" % nid}
    proj, enabled, lessons = bundle["proj"], bundle["enabled"], bundle["lessons"]
    has_weeks = any(l.get("week_no") for l in lessons)
    base, sid = _syno_login2(rows[0])
    try:
        dirs = [f for f in _syno_ls(base, sid, root) if f.get("isdir")]
        if not dirs:  # нПілНФ мЮРм≤ік∞А мЧЖк±∞лВШ(408) нХШмЬД нПілНФк∞А мЧЖмЭМ вАФ к≤љл°Ь мШ§л•Шл•Љ л™ЕнЩХнЮИ мХИлВі
            return {"ok": False, "error": "NASмЧРмДЬ к≥Љм†Х нПілНФл•Љ мЧі мИШ мЧЖмКµлЛИлЛ§ вАФ к≥Љм†ХмЭШ NAS нПілНФ к≤љл°Ьл•Љ нЩХмЭЄнХШмДЄмЪФ: nas%d:%s" % (nid, root)}
        st = next((s for s in enabled if s["id"] == LENGTH_STAGE_ID), None)
        want = (st.get("nas_folder") if st else "") or "мҐЕнОЄ"
        d7 = next((d for d in dirs if want in (d.get("name") or "") or "мҐЕнОЄ" in (d.get("name") or "")), None)
        if not d7:
            return {"ok": False, "error": "мҐЕнОЄ нПілНФл•Љ м∞ЊмІА л™їнЦИмКµлЛИлЛ§"}
        matches = {}  # basekey -> (mtime, path, name) вАФ м∞®мЛЬ лВі л™®лУ† нММнКЄ(мВђл≥Є¬ЈмИШм†Хл≥ЄмЭА мµЬмЛ† 1к∞Ь)
        for nm, pth, mt in _syno_walk(base, sid, d7["path"], 2):
            if os.path.splitext(nm)[1].lower() not in VIDEO_EXT:
                continue
            l = match_lesson(nm, lessons, has_weeks)
            if not l or l["id"] != lesson_id:
                continue
            k = _qc_basekey(nm)
            if k not in matches or mt >= matches[k][0]:
                matches[k] = (mt, pth, nm)
        if not matches:
            return {"ok": False, "error": "мЭі м∞®мЛЬмЧР лІ§мє≠лРШлКФ мҐЕнОЄ мШБмГБмЭД м∞ЊмІА л™їнЦИмКµлЛИлЛ§"}
        now = datetime.now(timezone.utc).isoformat()
        names, issues2 = [], []
        exist = {}
        try:  # м¶ЭлґД м†Рк≤А: нММмЭЉмЭі мХИ л∞ФлАМк≥† кЄ∞м§АлПД к∞ЩмЬЉл©і мЮђлґДмДЭ мГЭлЮµ (мЫРк≤©мЭА лЛ§мЪіл°ЬлУЬ лєДмЪ©мЭі мї§мДЬ нЪ®к≥Љ нБЉ)
            exist = {r0["file_name"]: r0 for r0 in (sb.table("audio_checks").select("file_name,file_mtime,qc_ver,status").eq("lesson_id", lesson_id).execute().data or [])}
        except Exception:
            pass
        for (mt, pth, nm) in sorted(matches.values(), key=lambda x: x[2]):
            e0 = exist.get(nm)
            if e0 and e0.get("file_mtime") == int(mt or 0) and e0.get("qc_ver") == QC_VER and e0.get("status") in ("ok", "warn", "bad"):
                names.append(nm)
                log("нТИмІИ м†Рк≤А(мГЭлЮµвАФл≥Ак≤љ мЧЖмЭМ, nas%s):" % nid, nm)
                continue
            row = {"lesson_id": lesson_id, "project_id": proj["id"], "file_name": nm,
                   "file_path": "nas%d:%s" % (nid, pth), "file_mtime": int(mt or 0), "qc_ver": QC_VER, "checked_at": now}
            try:
                issues, stats = _analyze_audio(_RemoteFSLite(base, sid), pth)
                bad = any(i["type"] in ("silence", "clip", "no_audio", "channel") for i in issues)
                warn = any(i["type"] in ("loud", "quiet", "jump", "spec", "black", "white", "contrast") for i in issues)
                row.update({"duration_sec": stats.get("duration"), "mean_volume": stats.get("mean"),
                            "max_volume": stats.get("max"), "issues": issues,
                            "status": "bad" if bad else ("warn" if warn else "ok"), "error": None})
                issues2.extend([dict(i, _file=nm) for i in issues])
            except Exception as e:
                row.update({"status": "error", "error": str(e), "issues": []})
            sb.table("audio_checks").upsert(row, on_conflict="lesson_id,file_name").execute()
            names.append(nm)
            log("нТИмІИ м†Рк≤А(nas%s):" % nid, nm, row.get("status"))
        try:  # мЭіл≤И м†Рк≤А лМАмГБмЧР мЧЖлКФ нММмЭЉмЭШ мШИм†Д к≤∞к≥Љ м†Хл¶ђ
            old = sb.table("audio_checks").select("id,file_name").eq("lesson_id", lesson_id).execute().data or []
            for o in old:
                if o.get("file_name") not in names:
                    sb.table("audio_checks").delete().eq("id", o["id"]).execute()
        except Exception:
            pass
        nm = ", ".join(names)
        if notify and issues2 and EMAIL_ENABLED:
            try:
                u = sb.table("users").select("email").eq("id", notify).execute().data
                email = u and u[0].get("email")
                if email:
                    TYPE_KR = {"silence": "лђімЭМ", "clip": "нБіл¶ђнХС", "loud": "к≥ЉлМАмЭМлЯЙ", "quiet": "к≥ЉмЖМмЭМлЯЙ",
                               "no_audio": "мШ§лФФмШ§ мЧЖмЭМ", "jump": "мЭМлЯЙ кЄЙл≥А", "channel": "м±ДлДР лђімЭМ"}
                    def mmss(v):
                        v = int(float(v or 0)); return "%d:%02d" % (v // 60, v % 60)
                    items = "".join("<li>[%s] %s ~ %s вАФ <b>%s</b> ¬Ј %s</li>" %
                                    (i.get("_file") or "", mmss(i.get("start")), mmss(i.get("end")),
                                     TYPE_KR.get(i.get("type"), i.get("type")), i.get("detail") or "")
                                    for i in issues2)
                    html = ("<div style=\"font-family:Malgun Gothic,sans-serif;font-size:14px;line-height:1.6\">"
                            "<p>мЧЕл°ЬлУЬнХШмЛ† <b>%s</b> мҐЕнОЄ мШБмГБмЧРмДЬ <b>нТИмІИ лђЄм†Ь %dк±і</b>мЭі л∞Ьк≤ђлРШмЧИмКµлЛИлЛ§.</p>"
                            "<p>нММмЭЉ: %s</p><ul>%s</ul>"
                            "<p>CDMSмЧРмДЬ нХілЛє м∞®мЛЬл•Љ нБіл¶≠нХШл©і 'рЯФК нТИмІИ м†Рк≤А' мДємЕШмЧР лђЄм†Ь кµђк∞ДмЭі нСЬмЛЬлР©лЛИлЛ§.</p></div>"
                            % (proj["name"], len(issues2), nm, items))
                    send_email(email, "[CDMS] нТИмІИ м†Рк≤А к≤∞к≥Љ вАФ %s лђЄм†Ь %dк±і" % (proj["name"], len(issues2)), html)
            except Exception as e:
                log("мШ§лФФмШ§ мХМл¶Љ л©ФмЭЉ мЛ§нМ®:", e)
        return {"ok": True, "checked": len(names), "problems": len(issues2)}
    finally:
        try:
            _syno_api(base, "/webapi/auth.cgi", {"api": "SYNO.API.Auth", "version": "6", "method": "logout",
                                                 "session": "FileStation", "_sid": sid})
        except Exception:
            pass


def action_scan_file(fs, params):
    import shutil
    p = params or {}
    path = p.get("path") or ""
    dest = p.get("dest") or ""
    name = p.get("name") or path.rsplit("/", 1)[-1]
    notify = p.get("notify_user")
    if not path or not dest:
        return {"ok": False, "error": "path/dest нХДмЪФ"}
    m = re.match(r"^nas(\d+):(.*)$", path)
    if m:  # лЛ§л•Є NAS(nasN:) вАФ FileStation APIл°Ь мЫРк≤© к≤АмВђ
        return _scan_file_remote(int(m.group(1)), m.group(2), re.sub(r"^nas\d+:", "", dest), name, notify)
    ok_exist = False
    for _ in range(20):  # мЧЕл°ЬлУЬ мІБнЫД нММмЭЉ мХИм∞© лМАкЄ∞ (мµЬлМА ~60міИ)
        if fs.exists(path):
            ok_exist = True
            break
        time.sleep(3)
    if not ok_exist:
        return {"ok": False, "error": "к≤АмВђ лМАмГБ нММмЭЉмЭД м∞ЊмІА л™їнЦИмКµлЛИлЛ§: %s" % path}
    scanner = shutil.which("clamdscan") or shutil.which("clamscan")
    if not scanner:
        return {"ok": False, "error": "мДЬл≤ДмЧР ClamAVк∞А мЧЖмКµлЛИлЛ§. ai-agentмЧРмДЬ: sudo apt-get install -y clamav clamav-daemon нЫД мЮђмЛЬлПД"}
    local, is_tmp = fs.local_copy(path)
    rc, sig, errtxt = 2, "", ""
    try:
        args = [scanner, "--no-summary"]
        if scanner.endswith("clamdscan"):
            args.append("--fdpass")
        out = subprocess.run(args + [local], capture_output=True, text=True, timeout=3600)
        rc = out.returncode
        errtxt = (out.stderr or "")[:200]
        for ln in (out.stdout or "").splitlines():
            if ": " in ln and not ln.rstrip().endswith("OK"):
                sig = ln.split(": ", 1)[1].strip()
                break
    finally:
        if is_tmp:
            try: os.unlink(local)
            except Exception: pass
    if rc == 0:  # м†ХмГБ вЖТ мµЬмҐЕ нПілНФл°Ь мЭілПЩ (к≤ємєШл©і мГИ мЭіл¶Д)
        base, ext = os.path.splitext(name)
        target = dest.rstrip("/") + "/" + name
        n = 1
        while fs.exists(target):
            n += 1
            target = dest.rstrip("/") + "/" + base + ("(%d)" % n) + ext
        fs.rename(path, target)
        log("л∞±мЛ†к≤АмВђ нЖµк≥Љ:", target)
        return {"ok": True, "clean": True, "moved": target}
    if rc == 1:  # мХЕмД±мљФлУЬ вЖТ к≤©л¶ђ + л©ФмЭЉ
        qdir = path.rsplit("/", 1)[0].rsplit("/", 1)[0] + "/.cdms_blocked"
        try: fs.makedirs(qdir)
        except Exception: pass
        qpath = qdir + "/" + name
        try: fs.rename(path, qpath)
        except Exception: qpath = path
        log("мХЕмД±мљФлУЬ м∞®лЛ®:", name, sig)
        if EMAIL_ENABLED:
            try:
                emails = []
                if notify:
                    u = sb.table("users").select("email").eq("id", notify).execute().data
                    if u and u[0].get("email"): emails.append(u[0]["email"])
                ad = sb.table("user_roles").select("user_id,users:user_id(email)").eq("role_code", "admin").execute().data
                for a in (ad or []):
                    e = (a.get("users") or {}).get("email")
                    if e: emails.append(e)
                html = ("<div style='font-family:Malgun Gothic,sans-serif;font-size:14px;line-height:1.6'>"
                        "<p>CDMSл°Ь мЧЕл°ЬлУЬлРЬ нММмЭЉмЧРмДЬ <b>мХЕмД±мљФлУЬк∞А нГРмІАлРШмЦі м∞®лЛ®</b>лРШмЧИмКµлЛИлЛ§.</p>"
                        "<p>нММмЭЉ: <b>%s</b><br>нГРмІАл™Е: %s<br>к≤©л¶ђ мЬДмєШ: %s</p>"
                        "<p>мЭі нММмЭЉмЭА NAS мЮСмЧЕ нПілНФмЧР л∞ШмШБлРШмІА мХКмХШмКµлЛИлЛ§. мЧЕл°ЬлУЬнХЬ PCмЭШ л∞±мЛ† м†Рк≤АмЭД кґМмЮ•нХ©лЛИлЛ§.</p></div>"
                        % (name, sig or "-", qpath))
                for em in sorted(set(emails)):
                    send_email(em, "[CDMS] вЫФ мЧЕл°ЬлУЬ м∞®лЛ® вАФ мХЕмД±мљФлУЬ нГРмІА: %s" % name, html)
            except Exception as e:
                log("м∞®лЛ® л©ФмЭЉ мЛ§нМ®:", e)
        return {"ok": True, "clean": False, "virus": sig or "malware", "quarantine": qpath}
    return {"ok": False, "error": "к≤АмВђ мШ§л•Ш(rc=%s): %s" % (rc, errtxt)}


# ============================================================================
# action: mkdir_tree / rename / sync_names / ping
# ============================================================================
def safe_name(s):
    return re.sub(r'[\\/:*?"<>|]+', "_", (s or "").strip()) or "untitled"


def action_mkdir_tree(fs, project, folders):
    root = project if (project or "").startswith("/") else safe_name(project)
    fs.makedirs(root)
    for f in folders or []:
        fs.makedirs("%s/%s" % (root, f))
    abs_root = root if os.path.isabs(root) else (
        os.path.join(NAS_BASE, root) if NAS_MODE == "mount" else r"\\%s\%s\%s" % (SMB_HOST, SMB_SHARE, root))
    return {"ok": True, "root": abs_root}


def action_rename_folder(fs, old, new):
    # new к∞А м†ИлМАк≤љл°Ьл©і кЈЄлМАл°Ь мВђмЪ©, мХДлЛИл©і old мЭШ лґАл™® мЬДмєШл•Љ мЬ†мІАнХШк≥† лІИмІАлІЙ segmentлІМ кµРм≤і
    if new and (new.startswith("/") or new.startswith("\\")):
        new_path = new
    else:
        nm = safe_name(new)
        sep = "\\" if (old and "\\" in old) else "/"
        parent = old[:old.rfind(sep)] if (old and sep in old) else ""
        new_path = (parent + sep + nm) if parent else nm
    if new_path != old and fs.exists(new_path):
        return {"ok": False, "error": "лМАмГБмЭі мЭілѓЄ м°імЮђнХі лНЃмЦімУ∞мІА мХКмХШмКµлЛИлЛ§: " + new_path}
    fs.rename(old, new_path)
    return {"ok": True, "old": old, "new": new_path}


def action_sync_names(fs, project_id):
    """NASмЭШ к≥Љм†Х л£®нКЄ нХШмЬД лЛ®к≥ДнПілНФл™ЕмЭД мЭљмЦі л∞ШнЩШ(к≤АнЖ†мЪ©). мЛ§м†Ь л∞ШмШБмЭА мЛ†м§СнЮИ."""
    b = load_project_bundle(project_id)
    if not b:
        return {"ok": False, "error": "project not found"}
    root = b["proj"].get("nas_root")
    if not root or not fs.exists(root):
        return {"ok": False, "error": "nas_root мЧЖмЭМ"}
    return {"ok": True, "folders": fs.listdirs(root)}


def action_ping(fs):
    return {"ok": True, "shares": fs.shares()}


# ============================================================================
# к∞ХмЭШк≥ДнЪНмДЬ нММмЛ± (м£Љм∞®¬Јм∞®мЛЬ мЮРлПЩ мґФмґЬ/мГЭмД±)
# ============================================================================
RE_WK = re.compile(r"(\d+)\s*м£Љ\s*м∞®")
RE_LS = re.compile(r"(\d+)\s*м∞®\s*мЛЬ")


def _syllabus_text(path, data):
    """мЧЕл°ЬлУЬлРЬ к∞ХмЭШк≥ДнЪНмДЬ л∞ФмЭінКЄ вЖТ нЕНмК§нКЄ (нШХмЛЭл≥Д)."""
    import tempfile, subprocess, zipfile
    ext = os.path.splitext(path)[1].lower()
    with tempfile.NamedTemporaryFile(delete=False, suffix=ext) as tf:
        tf.write(data); tmp = tf.name
    try:
        if ext == ".hwp":
            # hwp5txtлКФ нСЬ лВімЪ©мЭД "<нСЬ>"л°Ь к±ілДИлЫ∞лѓАл°Ь hwp5htmlл°Ь нСЬ(tr/td)кєМмІА мґФмґЬ
            try:
                import re as _re, html as _html
                hb = os.path.join(os.path.dirname(sys.executable), "hwp5html")
                cmd = hb if os.path.exists(hb) else "hwp5html"
                out = subprocess.run([cmd, "--html", tmp], capture_output=True, timeout=120).stdout.decode("utf-8", "ignore")
                lines = []
                for tr in _re.findall(r"<tr[^>]*>([\s\S]*?)</tr>", out):
                    cells = [_re.sub(r"\s+", " ", _html.unescape(_re.sub(r"<[^>]+>", " ", td))).strip()
                             for td in _re.findall(r"<td[^>]*>([\s\S]*?)</td>", tr)]
                    cells = [c for c in cells if c]
                    if not cells:
                        continue
                    lines.append("\t".join(cells))
                    # нСЬ нШХмЛЭ(м≤Ђ мЕА=мИЂмЮР, лЛ§мЭМ мЕА=м†Ьл™©)мЭД м∞®мЛЬ нМ®нДі лђЄмЮ•мЬЉл°ЬлПД нХ©мД± вЖТ _parse_syllabusк∞А мЭЄмЛЭ
                    if len(cells) >= 2 and _re.fullmatch(r"\d{1,2}", cells[0]):
                        title = cells[2] if (cells[1] in ("м£Љм†Ь", "м∞®мЛЬл™Е", "м£Љм∞®л™Е") and len(cells) >= 3) else cells[1]
                        if title and not _re.fullmatch(r"[\d\s.%]+", title):
                            lines.append("%sм∞®мЛЬ %s" % (cells[0], title))
                # л≥ЄлђЄ нЕНмК§нКЄлПД л≥СнЦЙ мґФмґЬ(нСЬ л∞Ц м∞®мЛЬ нСЬкЄ∞ лМАмЭС)
                txt = _re.sub(r"<[^>]+>", " ", _re.sub(r"</p>", "\n", out))
                lines.append(_html.unescape(txt))
                return "\n".join(lines)
            except Exception:
                return ""
        if ext == ".hwpx":
            out = []
            with zipfile.ZipFile(tmp) as z:
                for n in z.namelist():
                    if n.startswith("Contents/") and n.endswith(".xml"):
                        out.append(re.sub(r"<[^>]+>", " ", z.read(n).decode("utf-8", "ignore")))
            return "\n".join(out)
        if ext == ".pdf":
            import pdfplumber
            with pdfplumber.open(tmp) as pdf:
                return "\n".join((pg.extract_text() or "") for pg in pdf.pages)
        if ext == ".docx":
            import docx
            d = docx.Document(tmp); parts = [p.text for p in d.paragraphs]
            for t in d.tables:
                for row in t.rows:
                    parts.append("\t".join(c.text for c in row.cells))
            return "\n".join(parts)
        if ext == ".pptx":
            from pptx import Presentation
            prs = Presentation(tmp); parts = []
            for s in prs.slides:
                for sh in s.shapes:
                    if sh.has_text_frame:
                        for para in sh.text_frame.paragraphs:
                            parts.append("".join(r.text for r in para.runs))
                    if getattr(sh, "has_table", False):
                        for row in sh.table.rows:
                            parts.append("\t".join(c.text for c in row.cells))
            return "\n".join(parts)
        if ext in (".xlsx", ".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(tmp, read_only=True, data_only=True); parts = []
            for ws in wb.worksheets:
                for row in ws.iter_rows(values_only=True):
                    parts.append("\t".join("" if c is None else str(c) for c in row))
            return "\n".join(parts)
        if ext == ".csv":
            return data.decode("utf-8", "ignore")
        return data.decode("utf-8", "ignore")
    finally:
        try: os.unlink(tmp)
        except Exception: pass


def _parse_syllabus(text):
    lines = [re.sub(r"[ \t]+", " ", l).strip() for l in text.splitlines()]
    lines = [l for l in lines if l]
    weeks = {}; flat = []; cur = None
    strip = " :¬Ј.)¬Ј\t-"
    for ln in lines:
        wm = RE_WK.search(ln); lm = RE_LS.search(ln)
        if wm and not lm:
            wn = int(wm.group(1)); nm = RE_WK.sub("", ln).strip(strip)
            weeks.setdefault(wn, {"name": nm or None, "lessons": {}}); cur = wn; continue
        if lm:
            no = int(lm.group(1)); title = ln[lm.end():].strip(strip)
            if not title:
                title = RE_LS.sub("", ln).strip(strip)
            wmi = RE_WK.search(ln)
            if wmi:
                cur = int(wmi.group(1)); weeks.setdefault(cur, {"name": None, "lessons": {}})
            if cur is not None and weeks:
                weeks[cur]["lessons"][no] = title
            else:
                flat.append(title)
    return weeks, flat


def action_parse_syllabus(fs, project_id, path):
    if not (project_id and path):
        return {"ok": False, "error": "project_id/нММмЭЉк≤љл°Ь мЧЖмЭМ"}
    try:
        data = sb.storage.from_("plans").download(path)
    except Exception as e:
        return {"ok": False, "error": "мК§нЖ†л¶ђмІА лЛ§мЪіл°ЬлУЬ мЛ§нМ®: " + str(e)[:120]}
    text = _syllabus_text(path, data)
    if not text or len(text.strip()) < 5:
        return {"ok": False, "error": "нЕНмК§нКЄ мґФмґЬ мЛ§нМ®(нШХмЛЭ лѓЄмІАмЫР/лєИ нММмЭЉ). HWPлКФ hwp5txt(pyhwp) нХДмЪФ."}
    weeks, flat = _parse_syllabus(text)
    pst = sb.table("project_stages").select("stage_id").eq("project_id", project_id).eq("enabled", True).execute().data
    stage_ids = [r["stage_id"] for r in pst] or [1, 2, 3, 5, 7, 9]
    # кЄ∞м°і м£Љм∞®/м∞®мЛЬ мВ≠м†Ь нЫД мЮђмГЭмД±
    sb.table("lessons").delete().eq("project_id", project_id).execute()
    sb.table("weeks").delete().eq("project_id", project_id).execute()
    n = 0
    if weeks:
        for wn in sorted(weeks):
            wid = sb.table("weeks").insert({"project_id": project_id, "week_no": wn, "name": weeks[wn]["name"]}).execute().data[0]["id"]
            lessons = weeks[wn]["lessons"] or {1: None}
            for no in sorted(lessons):
                lid = sb.table("lessons").insert({"project_id": project_id, "week_id": wid, "lesson_no": no, "title": lessons[no]}).execute().data[0]["id"]
                sb.table("lesson_stage").insert([{"lesson_id": lid, "stage_id": s, "status": "wait"} for s in stage_ids]).execute()
                n += 1
        sb.table("projects").update({"course_type": "credit", "unit_label": "м£Љм∞®"}).eq("id", project_id).execute()
        return {"ok": True, "type": "м£Љм∞®нШХ", "weeks": len(weeks), "lessons": n}
    if not flat:
        return {"ok": False, "error": "м£Љм∞®/м∞®мЛЬ кµђм°∞л•Љ м∞ЊмІА л™їнЦИмКµлЛИлЛ§(нММмЭЉ лВімЪ© нЩХмЭЄ)."}
    for i, t in enumerate(flat, 1):
        lid = sb.table("lessons").insert({"project_id": project_id, "lesson_no": i, "title": t}).execute().data[0]["id"]
        sb.table("lesson_stage").insert([{"lesson_id": lid, "stage_id": s, "status": "wait"} for s in stage_ids]).execute()
        n += 1
    sb.table("projects").update({"course_type": "lesson", "unit_label": "м∞®мЛЬ"}).eq("id", project_id).execute()
    return {"ok": True, "type": "м∞®мЛЬнШХ", "lessons": n}


# ============================================================================
# мШБмГБк≤АмИШ: мҐЕнОЄ мШБмГБмЭД 480p нФДл°ЭмЛЬл°Ь л≥АнЩШ вЖТ Supabase мК§нЖ†л¶ђмІА мЧЕл°ЬлУЬ
# ============================================================================
def action_make_review_proxy(fs, project_id, lesson_id):
    import tempfile, subprocess
    if not (project_id and lesson_id):
        return {"ok": False, "error": "project_id/lesson_id мЧЖмЭМ"}
    b = load_project_bundle(project_id)
    if not b:
        return {"ok": False, "error": "project not found"}
    proj, lessons, enabled = b["proj"], b["lessons"], b["enabled"]
    root = proj.get("nas_root")
    if not root:
        return {"ok": False, "error": "nas_root лѓЄмД§м†Х"}
    lesson = next((l for l in lessons if l["id"] == lesson_id), None)
    if not lesson:
        return {"ok": False, "error": "lesson мЧЖмЭМ"}
    has_weeks = any(l.get("week_no") for l in lessons)
    st = next((s for s in enabled if s["id"] == LENGTH_STAGE_ID), None)
    folder = "%s/%s" % (root, (st or {}).get("nas_folder", "07_мҐЕнОЄ"))
    target = None
    for rel, mt in fs.walkfiles(folder):
        if os.path.splitext(rel)[1].lower() not in VIDEO_EXT:
            continue
        l = match_lesson(rel, lessons, has_weeks)
        if l and l["id"] == lesson_id:
            target = "%s/%s" % (folder, rel)
            break
    if not target:
        return {"ok": False, "error": "мЭі м∞®мЛЬмЭШ мҐЕнОЄ мШБмГБмЭД NASмЧРмДЬ м∞ЊмІА л™їнЦИмКµлЛИлЛ§."}
    # мЛЬлЖАл°ЬмІА/мЩЄлґА HTTPS мІБм†С мДЬлєЩ л™®лУЬ (л≥АнЩШ¬ЈмЧЕл°ЬлУЬ мЧЖмЭі к≥µк∞Ь URLлІМ кЄ∞л°Э)
    if NAS_PUBLIC_BASE:
        import urllib.parse
        absf = fs._abs(target) if hasattr(fs, "_abs") else target
        try:
            relpub = os.path.relpath(absf, NAS_PUBLIC_ROOT)
        except Exception:
            relpub = target
        url = NAS_PUBLIC_BASE.rstrip("/") + "/" + urllib.parse.quote(relpub.replace("\\", "/"))
        cur = sb.table("lessons").select("review_ver").eq("id", lesson_id).execute().data
        ver = ((cur[0].get("review_ver") if cur else 0) or 0) + 1
        sb.table("lessons").update({"review_path": url, "review_ver": ver}).eq("id", lesson_id).execute()
        return {"ok": True, "url": url, "version": ver, "mode": "synology"}
    local, is_tmp = fs.local_copy(target)
    out = tempfile.NamedTemporaryFile(delete=False, suffix=".mp4").name
    try:
        r = subprocess.run(["ffmpeg", "-y", "-i", local, "-vf", "scale=-2:480",
                            "-c:v", "libx264", "-crf", "30", "-preset", "veryfast",
                            "-c:a", "aac", "-b:a", "96k", "-movflags", "+faststart", out],
                           capture_output=True, timeout=900)
        if r.returncode != 0 or not os.path.exists(out) or os.path.getsize(out) == 0:
            return {"ok": False, "error": "ffmpeg л≥АнЩШ мЛ§нМ®"}
        cur = sb.table("lessons").select("review_ver").eq("id", lesson_id).execute().data
        ver = ((cur[0].get("review_ver") if cur else 0) or 0) + 1
        path = "%s/v%d.mp4" % (lesson_id, ver)
        with open(out, "rb") as f:
            data = f.read()
        try:
            sb.storage.from_("review").upload(path, data, {"content-type": "video/mp4", "upsert": "false"})  # лНЃмЦімУ∞кЄ∞ кЄИмІА
        except Exception as e:
            return {"ok": False, "error": "мЧЕл°ЬлУЬ мЛ§нМ®: " + str(e)[:140]}
        sb.table("lessons").update({"review_path": path, "review_ver": ver}).eq("id", lesson_id).execute()
        return {"ok": True, "path": path, "version": ver, "size_mb": round(len(data) / 1e6, 1)}
    finally:
        for p in (out,):
            try: os.unlink(p)
            except Exception: pass
        if is_tmp:
            try: os.unlink(local)
            except Exception: pass


# ============================================================================
# нБР м≤Шл¶ђ
# ============================================================================
def action_make_preview(fs, project_id, path):
    """PSD¬ЈмЭілѓЄмІА лѓЄл¶ђл≥ікЄ∞ мГЭмД± вАФ лЛ§мЪіл°ЬлУЬ мЧЖмЭі лЄМлЭЉмЪ∞м†АмЧРмДЬ нЩХмЭЄ (previews л≤ДнВЈмЧР JPG мЇРмЛЬ).

    PSDлКФ нММмЭЉмЧР лВімЮ•лРЬ л≥СнХ© мЭілѓЄмІА(composite)л•Љ Pillowл°Ь мЭљмЦі 1600px JPGл°Ь л≥АнЩШнХЬлЛ§.
    нПђнЖ†мГµ м†АмЮ• мЛЬ 'нШЄнЩШмД± мµЬлМАнЩФ'к∞А кЇЉмІД PSDлКФ л≥СнХ© мЭілѓЄмІАк∞А мЧЖмЦі мЛ§нМ®нХ† мИШ мЮИлЛ§.
    """
    import hashlib
    try:
        from PIL import Image
    except Exception:
        return {"ok": False, "error": "мДЬл≤ДмЧР Pillowк∞А мЧЖмКµлЛИлЛ§ (pip install pillow)"}
    path = str(path or "")
    if not path:
        return {"ok": False, "error": "path нХДмЪФ"}
    ext = os.path.splitext(path)[1].lower()
    if ext not in (".psd", ".psb", ".jpg", ".jpeg", ".png", ".gif", ".bmp", ".webp", ".tif", ".tiff"):
        return {"ok": False, "error": "лѓЄл¶ђл≥ікЄ∞ лѓЄмІАмЫР нШХмЛЭ: %s" % ext}
    key = hashlib.sha256(path.encode("utf-8")).hexdigest()  # нФДлЯ∞нКЄмЩА лПЩмЭЉ кЈЬмєЩ(SubtleCrypto SHA-256)
    dest = "psd/%s.jpg" % key
    pub = "%s/storage/v1/object/public/previews/%s" % (SB_URL, dest)
    try:  # мЇРмЛЬ вАФ к∞ЩмЭА к≤љл°ЬлКФ мЮђмГЭмД±нХШмІА мХКмЭМ
        ex = sb.storage.from_("previews").list("psd", {"search": key})
        if any((it.get("name") == "%s.jpg" % key) for it in (ex or [])):
            return {"ok": True, "url": pub, "cached": True}
    except Exception:
        pass
    m = re.match(r"^nas(\d+):(.*)$", path)
    local, is_tmp = None, False
    try:
        if m:  # лЛ§л•Є NAS(nasN:) вАФ FileStation APIл°Ь лВіл†§л∞ЫмХД л≥АнЩШ
            rows = sb.table("nas_config").select("*").eq("id", int(m.group(1))).execute().data
            if not rows:
                return {"ok": False, "error": "nas_config(id=%s) мЧЖмЭМ" % m.group(1)}
            base, sid = _syno_login2(rows[0])
            tf = tempfile.NamedTemporaryFile(prefix="cdms_pv_", delete=False)
            local, is_tmp = tf.name, True
            tf.close()
            _syno_download(base, sid, m.group(2), local)
        else:
            local, is_tmp = fs.local_copy(path)
        im = Image.open(local)
        im.load()
        if im.mode not in ("RGB", "L"):
            im = im.convert("RGB")
        im.thumbnail((1600, 1600))
        out = tempfile.NamedTemporaryFile(prefix="cdms_pv_out_", suffix=".jpg", delete=False)
        outp = out.name
        out.close()
        im.save(outp, "JPEG", quality=82)
        with open(outp, "rb") as fjpg:
            data = fjpg.read()
        try:
            os.unlink(outp)
        except Exception:
            pass
    except Exception as e:
        return {"ok": False, "error": "мЭілѓЄмІА нХімДЭ мЛ§нМ® вАФ нПђнЖ†мГµ 'нШЄнЩШмД± мµЬлМАнЩФ' мЧЖмЭі м†АмЮ•лРЬ PSDмЭЉ мИШ мЮИмКµлЛИлЛ§: %s" % e}
    finally:
        if is_tmp and local:
            try:
                os.unlink(local)
            except Exception:
                pass
    try:
        sb.storage.from_("previews").upload(dest, data, {"content-type": "image/jpeg", "upsert": "true"})
    except Exception as e:
        es = str(e).lower()
        if "exist" not in es and "duplicate" not in es:
            return {"ok": False, "error": "лѓЄл¶ђл≥ікЄ∞ мЧЕл°ЬлУЬ мЛ§нМ®: %s" % e}
    return {"ok": True, "url": pub}


def dispatch(fs, task):
    action = task["action"]
    p = task.get("params") or {}
    pid = task.get("project_id") or p.get("project_id")
    if action == "ping":
        return action_ping(fs)
    if action == "mkdir_tree":
        return action_mkdir_tree(fs, p.get("project"), p.get("folders"))
    if action == "scan_progress":
        return action_scan_progress(fs, pid)
    if action == "probe_durations":
        return action_probe_durations(fs, pid)
    if action == "audio_check":
        return action_audio_check(fs, pid, p.get("lesson_id"), p.get("notify_user"))
    if action == "scan_file":
        return action_scan_file(fs, p)
    if action == "rename_folder":
        return action_rename_folder(fs, p.get("old"), p.get("new"))
    if action == "sync_names":
        return action_sync_names(fs, pid)
    if action == "parse_syllabus":
        return action_parse_syllabus(fs, pid, p.get("path"))
    if action == "make_review_proxy":
        return action_make_review_proxy(fs, pid, p.get("lesson_id"))
    if action == "make_preview":
        return action_make_preview(fs, pid, p.get("path"))
    return {"ok": False, "error": "unknown action: %s" % action}


def process_queue(fs):
    stale = (datetime.now(timezone.utc) - timedelta(hours=3)).isoformat()
    rows = sb.table("nas_tasks").select("*").or_(
        "status.is.null,status.eq.pending,and(status.eq.running,updated_at.lt.%s)" % stale).order(
        "created_at").limit(10).execute().data
    for t in rows:
        log("мЮСмЧЕ м≤Шл¶ђ:", t["action"], t.get("project_id") or "")
        try:  # м≤Шл¶ђ мЛЬмЮС нСЬмЛЬ вАФ нФДлЯ∞нКЄмЧРмДЬ лМАкЄ∞/мІДнЦЙ кµђлґД, м§СлЛ® мЛЬ 3мЛЬк∞Д нЫД мЮРлПЩ мЮђмЛЬлПД
            sb.table("nas_tasks").update({"status": "running", "updated_at": datetime.now(timezone.utc).isoformat()}).eq("id", t["id"]).execute()
        except Exception:
            pass
        try:
            res = dispatch(fs, t)
        except Exception as e:
            res = {"ok": False, "error": str(e)}
            traceback.print_exc()
        status = "done" if res.get("ok") else "error"
        sb.table("nas_tasks").update(
            {"status": status, "result": res,
             "updated_at": datetime.now(timezone.utc).isoformat()}).eq("id", t["id"]).execute()


def auto_scan(fs):
    projs = sb.table("projects").select("id,name,nas_root").execute().data
    for p in projs:
        if not p.get("nas_root"):
            continue
        try:
            r = action_scan_progress(fs, p["id"])
            log("мЮРлПЩмК§мЇФ:", p["name"], r)
        except Exception as e:
            log("мЮРлПЩмК§мЇФ мЛ§нМ®:", p["name"], e)


def check_contract_reminders():
    """к≥ДмХљ мЩДл£МмЭЉ 30/14/7мЭЉ м†Д вЖТ PM(лШРлКФ нПіл∞± кіАл¶ђмЮР)мЧРк≤М мХИлВіл©ФмЭЉ. нХШл£® 1нЪМ."""
    if not EMAIL_ENABLED:
        return 0
    from datetime import date
    today = date.today()
    try:
        progs = sb.table("programs").select("id,name,contract_end,pm_id").not_.is_("contract_end", "null").execute().data
    except Exception as e:
        log("к≥ДмХљлІМл£М м†Рк≤А мЛ§нМ®:", e); return 0
    users = {u["id"]: u for u in sb.table("users").select("id,name,email").execute().data}
    done = set((r["ref"], r["label"]) for r in sb.table("reminders").select("ref,label").eq("kind", "contract_expiry").execute().data)
    sent = 0
    for p in progs:
        try:
            ce = datetime.fromisoformat(str(p["contract_end"])).date()
        except Exception:
            continue
        days = (ce - today).days
        for th, label in ((30, "d30"), (14, "d14"), (7, "d7")):
            if days != th or (p["id"], label) in done:
                continue
            u = users.get(p.get("pm_id")) if p.get("pm_id") else None
            to = (u or {}).get("email") or ""
            if (not to) or ("mirim.local" in to) or ("mirimcms.local" in to):
                to = REMIND_EMAIL
            subj = "[CDMS] к≥ДмХљ лІМл£М D-%d ¬Ј %s" % (th, p["name"])
            html = ("<div style='font-family:sans-serif;font-size:14px'><p><b>%s</b> мВђмЧЕмЭШ к≥ДмХљ мЩДл£МмЭЉмЭі "
                    "<b>%s</b> мЮЕлЛИлЛ§ (D-%d).</p><p>мІДнЦЙ нШДнЩ©: <a href='%s'>%s</a></p>"
                    "<p style='color:#888;font-size:12px'>вАФ CDMS мЮРлПЩмХМл¶Љ</p></div>") % (
                    p["name"], ce.isoformat(), th, CDMS_URL, CDMS_URL)
            ok, err = send_email(to, subj, html)
            try:
                sb.table("reminders").insert({"kind": "contract_expiry", "ref": p["id"], "label": label,
                                              "to_email": to, "status": "sent" if ok else "error", "error": err}).execute()
            except Exception as e:
                log("reminders кЄ∞л°Э мЛ§нМ®:", e)
            if ok:
                sent += 1
            else:
                log("к≥ДмХљлІМл£М л©ФмЭЉ мЛ§нМ®:", p["name"], to, err)
    return sent


def selfcheck():
    """мД§мєШ мЛЬ мЧ∞к≤∞ мЮРк∞АмІДлЛ®. л™®лСР нЖµк≥ЉнХШл©і exit 0, нХШлВШлЭЉлПД мЛ§нМ®нХШл©і exit 1."""
    import shutil
    ok = True
    # 1) ffprobe
    if shutil.which("ffprobe"):
        print("  [PASS] ffprobe мД§мєШлР®")
    else:
        print("  [FAIL] ffprobe мЧЖмЭМ вАФ sudo apt-get install -y ffmpeg"); ok = False
    # 2) Supabase (service_role)
    try:
        n = sb.table("stages").select("id", count="exact").limit(1).execute()
        print("  [PASS] Supabase мЧ∞к≤∞ OK (stages м†СкЈЉ)")
    except Exception as e:
        print("  [FAIL] Supabase мЧ∞к≤∞ мЛ§нМ®:", str(e)[:160]); ok = False
    # 3) NAS
    try:
        fs = build_fs()
        if NAS_MODE == "smb":
            sh = fs.shares()
            print("  [PASS] NAS(SMB) м†СкЈЉ OK:", (sh[0] if sh else ""))
        else:
            if os.path.isdir(NAS_BASE):
                print("  [PASS] NAS лІИмЪінКЄ к≤љл°Ь OK:", NAS_BASE)
            else:
                print("  [FAIL] NAS_BASE к≤љл°Ь мЧЖмЭМ:", NAS_BASE, "(лІИмЪінКЄ нЩХмЭЄ)"); ok = False
    except Exception as e:
        print("  [FAIL] NAS м†СкЈЉ мЛ§нМ®:", str(e)[:160]); ok = False
    # 4) мЭіл©ФмЭЉ
    if EMAIL_ENABLED:
        if EMAIL_PROVIDER == "smtp":
            if SMTP_USER and SMTP_PASS:
                print("  [PASS] мЭіл©ФмЭЉ SMTP мД§м†ХлР® (%s:%d, %s)" % (SMTP_HOST, SMTP_PORT, SMTP_USER))
            else:
                print("  [WARN] EMAIL_PROVIDER=smtp мЭЄлН∞ SMTP_USER/SMTP_PASS лєДмЦімЮИмЭМ вАФ л©ФмЭЉ лѓЄл∞ЬмЖ°")
        elif EMAIL_API_KEY:
            print("  [PASS] мЭіл©ФмЭЉ нВ§ мД§м†ХлР® (provider=%s)" % EMAIL_PROVIDER)
        else:
            print("  [WARN] EMAIL_ENABLED=true мЭЄлН∞ EMAIL_API_KEY лєДмЦімЮИмЭМ вАФ л©ФмЭЉ лѓЄл∞ЬмЖ°")
    else:
        print("  [INFO] мЭіл©ФмЭЉ лєДнЩЬмД±(EMAIL_ENABLED=false)")
    return ok


def main():
    log("CDMS NAS Worker мЛЬмЮС вАФ mode=%s base=%s мҐЕнОЄлЛ®к≥Д=%d" % (NAS_MODE, NAS_BASE, LENGTH_STAGE_ID))
    fs = build_fs()
    last_scan = 0
    last_remind = None
    while True:
        try:
            process_queue(fs)
            if SCAN_INTERVAL > 0 and (time.time() - last_scan) >= SCAN_INTERVAL:
                auto_scan(fs)
                last_scan = time.time()
            today = datetime.now().strftime("%Y-%m-%d")
            if last_remind != today:
                n = check_contract_reminders()
                if n:
                    log("к≥ДмХљлІМл£М мХМл¶Љ л∞ЬмЖ°:", n)
                last_remind = today
        except Exception as e:
            log("л£®нФД мШ§л•Ш:", e)
            traceback.print_exc()
        time.sleep(POLL_INTERVAL)


if __name__ == "__main__":
    cmd = sys.argv[1] if len(sys.argv) > 1 else ""
    if cmd == "selfcheck":
        print("== CDMS NAS Worker мЮРк∞АмІДлЛ® ==")
        sys.exit(0 if selfcheck() else 1)
    elif cmd == "scan-once":
        fs = build_fs()
        target = sys.argv[2] if len(sys.argv) > 2 else None
        if target:
            print(json.dumps(action_scan_progress(fs, target), ensure_ascii=False, indent=2))
        else:
            auto_scan(fs)
    else:
        main()
